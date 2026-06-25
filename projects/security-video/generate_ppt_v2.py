#!/usr/bin/env python3
"""
上海银行安防宣传片 PPT v2 — 深色底 + 大字卡片 + 演讲型
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import openpyxl
import os, re
from pathlib import Path

EXCEL_PATH = "/root/.openclaw/media/qqbot/downloads/20260616上海银行安防数字化探索实践短片脚本-V1_1781602866070_644fd5.xlsx"
OUTPUT_DIR = Path("/root/.openclaw/workspace/projects/security-video/output")
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# ===== 配色 =====
BG_DARK   = RGBColor(0x0B, 0x0E, 0x14)   # 深黑蓝底
BG_CARD   = RGBColor(0x15, 0x1A, 0x24)   # 卡片背景
TECH_BLUE = RGBColor(0x38, 0x7F, 0xF7)   # 科技蓝
ACCENT_CYAN = RGBColor(0x22, 0xD3, 0xEE) # 青蓝色点缀
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB8, 0xC4)
MID_GRAY  = RGBColor(0x70, 0x78, 0x88)
GOLD      = RGBColor(0xF5, 0x9E, 0x0B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

def parse_scenes():
    wb = openpyxl.load_workbook(EXCEL_PATH)
    ws = wb['上篇']
    scenes = []
    current_chapter = ""
    current_section = ""
    for row in ws.iter_rows(values_only=True):
        vals = [str(c).strip() if c else '' for c in row[:6]]
        first = vals[0]
        if all(v == '' for v in vals[:4]): continue
        if '章节' in first: current_chapter = first; continue
        if '板块' in first: current_section = first; continue
        if '部分' in first or '镜号' in first or '此篇' in first: continue
        if first.isdigit():
            scenes.append({
                'num': int(first), 'time': vals[1],
                'visual': vals[2].replace('<br>', '\n'),
                'subtitle': vals[3].replace('<br>', '\n'),
                'narration': vals[4].replace('<br>', '\n'),
                'materials': vals[5] if len(vals) > 5 else '',
                'chapter': current_chapter, 'section': current_section,
            })
    return scenes

def add_bg(slide, left=0, top=0, w=None, h=None, color=BG_DARK):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(left) if isinstance(left,(int,float)) else left,
        Inches(top) if isinstance(top,(int,float)) else top,
        SLIDE_W if w is None else (Inches(w) if isinstance(w,(int,float)) else w),
        SLIDE_H if h is None else (Inches(h) if isinstance(h,(int,float)) else h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def add_card(slide, left, top, w, h, color=BG_CARD, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
    return shape

def add_text(slide, left, top, w, h, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font="微软雅黑"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color; p.font.name = font
    p.alignment = align
    return box

def add_lines(slide, left, top, w, h, lines, align=PP_ALIGN.LEFT):
    """lines: list of (text, size, bold, color)"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.font.size = Pt(size); p.font.bold = bold
        p.font.color.rgb = color; p.font.name = "微软雅黑"; p.alignment = align
        p.space_after = Pt(4)
    return box

def add_accent_bar(slide, left, top, w=0.06, h=0.6, color=TECH_BLUE):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def add_line_decor(slide, left, top, w, color=TECH_BLUE, h=0.015):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def extract_key_points(subtitle_text, max_points=5):
    """从字幕文本提取核心要点"""
    text = subtitle_text.strip()
    if not text or text == '/':
        return []
    
    # 清理标记
    text = text.replace('<br>', '\n')
    
    points = []
    for line in text.split('\n'):
        line = line.strip()
        if not line: continue
        
        # 去掉【下阶段】
        if '【下阶段】' in line:
            continue
        
        # 数字开头 → 条目
        m = re.match(r'^(\d+)[、.]\s*(.+)', line)
        if m:
            points.append(m.group(2).strip())
            continue
        
        # 其他行
        if len(line) > 5:
            # 去段落标记
            line = re.sub(r'^[▸►▪·•\-\*]\s*', '', line)
            points.append(line)
    
    # 去重，限制条数
    seen = set()
    result = []
    for p in points:
        if p not in seen:
            seen.add(p)
            result.append(p)
    
    return result[:max_points]


def make_cover(prs):
    """封面 — 深色大标题"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    
    # 顶部光条
    add_line_decor(slide, 0, 0, 13.333, TECH_BLUE, 0.04)
    
    # 左侧渐变装饰
    add_bg(slide, 0, 0, 0.3, 7.5, BG_CARD)
    
    # 标题
    add_text(slide, 2.0, 1.8, 9.5, 1.5,
             "上海银行", 56, WHITE, True)
    add_text(slide, 2.0, 2.6, 9.5, 1.0,
             "安防数智化应用展示", 44, TECH_BLUE, True)
    
    # 装饰线
    add_line_decor(slide, 2.0, 3.7, 3.0)
    
    # 副标题
    add_text(slide, 2.0, 4.0, 9.5, 0.8,
             "统筹发展与安全，筑牢金融安全防线", 24, LIGHT_GRAY, False)
    add_text(slide, 2.0, 4.7, 9.5, 0.5,
             "上篇 · 安防数智化探索实践", 16, MID_GRAY, False)


def make_chapter(slide, title, subtitle, tag=""):
    """章节分隔页"""
    add_bg(slide)
    add_line_decor(slide, 0, 0, 13.333, TECH_BLUE, 0.04)
    
    y = 2.2
    if tag:
        add_text(slide, 1.5, y, 10.5, 0.5, tag, 14, TECH_BLUE, True)
        y += 0.6
    
    add_text(slide, 1.5, y, 10.5, 1.2, title, 48, WHITE, True)
    
    if subtitle:
        add_line_decor(slide, 1.5, y + 1.4, 2.5)
        add_text(slide, 1.5, y + 1.7, 10.5, 1.0, subtitle, 22, LIGHT_GRAY, False)
    
    # 右下角页码装饰
    add_bg(slide, 0, 6.5, 13.333, 1.0, BG_CARD)


def make_content(prs, page_num, total, scene, section_label=""):
    """内容页 — 深色底 + 标题 + 卡片要点"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_line_decor(slide, 0, 0, 13.333, TECH_BLUE, 0.03)
    
    # === 顶部信息条 ===
    add_text(slide, 0.8, 0.25, 2.0, 0.4,
             f"0{page_num}" if page_num < 10 else str(page_num),
             26, TECH_BLUE, True, PP_ALIGN.LEFT, "Segoe UI")
    
    if section_label:
        add_text(slide, 2.5, 0.3, 8.0, 0.4, section_label, 12, MID_GRAY, False)
    
    add_text(slide, 11.0, 0.3, 2.0, 0.4,
             f"{page_num}/{total}", 10, MID_GRAY, False, PP_ALIGN.RIGHT, "Segoe UI")
    
    add_line_decor(slide, 0.8, 0.75, 11.733, RGBColor(0x25, 0x2D, 0x3D), 0.008)
    
    # === 标题（大字） ===
    subtitle = scene['subtitle'].strip()
    title_text = ""
    key_points = []
    
    if subtitle and subtitle != '/':
        lines = subtitle.split('\n')
        # 第一行通常是最重要的标题
        first = lines[0].strip()
        first = first.replace('【下阶段】', '')
        title_text = first
        key_points = extract_key_points(subtitle, 5)
    else:
        title_text = f"镜号 {scene['num']} | {scene['time']}"
        key_points = []
    
    # 没有提取到要点时，用旁白提炼
    if not key_points:
        nar = scene['narration'].strip()
        if nar and nar != '/':
            key_points = extract_key_points(nar, 4)
    
    # === 左列：标题 ===
    y_start = 1.3
    add_text(slide, 0.8, y_start, 11.7, 0.9, title_text, 32, WHITE, True)
    add_line_decor(slide, 0.8, y_start + 1.0, 2.0)
    
    # === 核心要点卡片 ===
    if key_points:
        card_y = y_start + 1.5
        # 1-2个要点用大卡片
        if len(key_points) <= 2:
            for i, pt in enumerate(key_points):
                cy = card_y + i * 2.3
                add_card(slide, 0.8, cy, 11.733, 1.8)
                # 序号
                add_text(slide, 1.2, cy + 0.3, 0.6, 0.6,
                         f"0{i+1}", 28, TECH_BLUE, True, PP_ALIGN.LEFT, "Segoe UI")
                # 要点文字
                add_text(slide, 2.0, cy + 0.35, 10.0, 1.3,
                         pt, 20, WHITE, False)
        else:
            # 3-5个要点用横排卡片
            n = len(key_points)
            card_w = (11.733 - 0.2 * (n - 1)) / n
            for i, pt in enumerate(key_points):
                cx = 0.8 + i * (card_w + 0.2)
                add_card(slide, cx, card_y, card_w, 2.5)
                add_text(slide, cx + 0.3, card_y + 0.2, card_w - 0.6, 0.5,
                         f"0{i+1}", 24, TECH_BLUE, True, PP_ALIGN.LEFT, "Segoe UI")
                # 截取要点
                if len(pt) > 60:
                    pt = pt[:60] + "..."
                add_text(slide, cx + 0.3, card_y + 0.8, card_w - 0.6, 1.5,
                         pt, 16, WHITE, False)
    
    # === 底部 :: 旁白核心句 ===
    nar = scene['narration'].strip()
    if nar and nar != '/':
        # 取第一句或核心句
        nar_short = nar.split('。')[0] if '。' in nar else nar[:120]
        if len(nar_short) > 120:
            nar_short = nar_short[:120] + "..."
        
        # 底部区域
        add_bg(slide, 0.8, 6.2, 11.733, 1.0, BG_CARD)
        add_text(slide, 1.0, 6.3, 0.8, 0.3, "旁  白", 10, TECH_BLUE, True)
        add_text(slide, 1.0, 6.55, 11.3, 0.55, nar_short, 12, LIGHT_GRAY, False)
    
    # 素材标记
    mats = scene.get('materials', '').strip()
    if mats and mats != '/':
        add_text(slide, 11.5, 6.8, 1.5, 0.25, "📎 素材", 8, MID_GRAY, False, PP_ALIGN.RIGHT)


def make_ending(prs):
    """结尾页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_line_decor(slide, 0, 0, 13.333, TECH_BLUE, 0.04)
    
    add_text(slide, 1.5, 2.0, 10.5, 1.0,
             "持续数智创新  筑牢金融安全屏障", 40, WHITE, True, PP_ALIGN.CENTER)
    add_line_decor(slide, 5.5, 3.1, 2.5)
    add_text(slide, 1.5, 3.4, 10.5, 0.8,
             "金融让生活更美好", 24, TECH_BLUE, False, PP_ALIGN.CENTER)
    add_text(slide, 1.5, 4.5, 10.5, 0.5,
             "上海银行  |  安全保卫部", 16, MID_GRAY, False, PP_ALIGN.CENTER)


def main():
    scenes = parse_scenes()
    print(f"📋 {len(scenes)} 个场景")
    
    prs = Presentation()
    prs.slide_width = Emu(int(SLIDE_W))
    prs.slide_height = Emu(int(SLIDE_H))
    
    p = 0  # page counter
    total_hint = 20
    
    # 封面
    print("封面")
    make_cover(prs); p += 1
    
    # === Ch1: 开篇 ===
    print("Ch1 开篇")
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    make_chapter(s1, "开篇 · 背景与发展概况",
                 "模拟安防 → 数字安防 → 智能应用 → 数智化阶段", "CHAPTER 01")
    p += 1
    for s in scenes:
        if 1 <= s['num'] <= 4:
            make_content(prs, p, total_hint, s, "第一章 开篇与背景")
            p += 1
    
    # === Ch2: 规划框架 ===
    print("Ch2 规划")
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    make_chapter(s2, "整体规划框架",
                 '三大方向 · "看得到 管得牢 用得优"', "CHAPTER 02")
    p += 1
    for s in scenes:
        if s['num'] == 5:
            make_content(prs, p, total_hint, s, "第二章 规划框架")
            p += 1
    
    # === Ch3: 核心成果 ===
    print("Ch3 核心成果")
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    make_chapter(s3, "核心实践成果",
                 '三大板块 · "看管用"全面落地', "CHAPTER 03")
    p += 1
    
    for s in scenes:
        if 6 <= s['num'] <= 13:
            sec = s.get('section', '') or ''
            label = "第三章 核心成果"
            if '看得到' in sec: label = "板块一 视觉与监测突破"
            elif '防得住' in sec: label = "板块二 全流程管理创新"
            elif '用得优' in sec: label = "板块三 AI赋能业务与基层"
            make_content(prs, p, total_hint, s, label)
            p += 1
    
    # === Ch4: 总结 ===
    print("Ch4 总结")
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    make_chapter(s4, "上篇 · 总结展望",
                 "立足安防 · 拥抱AI · 服务全行", "CHAPTER 04")
    p += 1
    for s in scenes:
        if s['num'] == 15:
            make_content(prs, p, total_hint, s, "第四章 总结展望")
            p += 1
    
    # 结尾
    print("结尾")
    make_ending(prs); p += 1
    
    out = str(OUTPUT_DIR / "上海银行安防宣传片_PPT_v2.pptx")
    prs.save(out)
    print(f"\n✅ {out}  |  {os.path.getsize(out)/1024:.0f} KB  |  {p} 页")


if __name__ == '__main__':
    main()

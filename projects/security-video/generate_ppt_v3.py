#!/usr/bin/env python3
"""
上海银行安防宣传片 PPT v3 — 大字少留白，视频背景风格
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import openpyxl
import os, re
from pathlib import Path

EXCEL_PATH = "/root/.openclaw/media/qqbot/downloads/20260616上海银行安防数字化探索实践短片脚本-V1_1781602866070_644fd5.xlsx"
OUTPUT_DIR = Path("/root/.openclaw/workspace/projects/security-video/output")
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# 配色
BG_DARK    = RGBColor(0x08, 0x0C, 0x14)
BG_CARD    = RGBColor(0x12, 0x18, 0x26)
TECH_BLUE  = RGBColor(0x3B, 0x82, 0xF6)
LIGHT_BLUE = RGBColor(0x93, 0xC5, 0xFD)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0x9C, 0xA3, 0xAF)
DARK_GRAY  = RGBColor(0x6B, 0x72, 0x80)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN = 0.6

def parse_scenes():
    wb = openpyxl.load_workbook(EXCEL_PATH)
    ws = wb['上篇']
    scenes = []
    current_chapter = ""
    current_section = ""
    for row in ws.iter_rows(values_only=True):
        vals = [str(c).strip() if c else '' for c in row[:6]]
        first = vals[0]
        if all(v == '' for v in vals[:4]): continue
        if '章节' in first: current_chapter = first; continue
        if '板块' in first: current_section = first; continue
        if '部分' in first or '镜号' in first or '此篇' in first: continue
        if first.isdigit():
            scenes.append({
                'num': int(first), 'time': vals[1],
                'visual': vals[2].replace('<br>', '\n'),
                'subtitle': vals[3].replace('<br>', '\n'),
                'narration': vals[4].replace('<br>', '\n'),
                'materials': vals[5] if len(vals) > 5 else '',
                'chapter': current_chapter, 'section': current_section,
            })
    return scenes

def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

def card(slide, l, t, w, h, color=BG_CARD):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

def txt(slide, l, t, w, h, text, size=22, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font="微软雅黑"):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color; p.font.name = font
    p.alignment = align
    return box

def add_lines(slide, l, t, w, h, lines, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.font.size = Pt(size); p.font.bold = bold
        p.font.color.rgb = color; p.font.name = "微软雅黑"; p.alignment = align
        p.space_after = Pt(6)
    return box

def bar(slide, l, t, w, h=0.02, color=TECH_BLUE):
    rect(slide, l, t, w, h, color)

def extract_key_points(text, max_pts=4):
    if not text or text == '/': return []
    pts = []
    for line in text.replace('<br>', '\n').split('\n'):
        line = line.strip()
        if not line or '【下阶段】' in line: continue
        line = re.sub(r'^[▸►▪·•\-\*]\s*', '', line)
        m = re.match(r'^\d+[、.]\s*(.+)', line)
        if m: line = m.group(1)
        if len(line) > 4:
            pts.append(line)
    seen = set(); result = []
    for p in pts:
        if p not in seen: seen.add(p); result.append(p)
    return result[:max_pts]


def make_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, 13.333, 7.5, BG_DARK)
    bar(slide, 0, 0, 13.333, 0.05)
    
    # 大字标题 — 撑满画面
    txt(slide, 0.6, 1.5, 12.1, 1.6, "上海银行", 64, WHITE, True)
    txt(slide, 0.6, 2.7, 12.1, 1.2, "安防数智化应用展示", 50, TECH_BLUE, True)
    bar(slide, 0.6, 4.0, 4.0)
    txt(slide, 0.6, 4.3, 12.1, 0.8, "统筹发展与安全，筑牢金融安全防线", 28, GRAY)
    txt(slide, 0.6, 5.2, 12.1, 0.5, "上篇 · 安防数智化探索实践", 20, DARK_GRAY)


def make_chapter(slide, chapter_num, title, subtitle=""):
    rect(slide, 0, 0, 13.333, 7.5, BG_DARK)
    bar(slide, 0, 0, 13.333, 0.05)
    
    txt(slide, 0.6, 0.8, 4.0, 0.6, f"CHAPTER {chapter_num:02d}", 20, TECH_BLUE, True, PP_ALIGN.LEFT, "Segoe UI")
    
    txt(slide, 0.6, 2.2, 12.1, 1.5, title, 52, WHITE, True)
    
    if subtitle:
        bar(slide, 0.6, 3.8, 3.0)
        txt(slide, 0.6, 4.1, 12.1, 1.0, subtitle, 26, GRAY)
    
    # 右下装饰
    rect(slide, 0, 6.0, 13.333, 1.5, BG_CARD)


def make_content(prs, page_num, total, scene, section_label=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, 13.333, 7.5, BG_DARK)
    bar(slide, 0, 0, 13.333, 0.04)
    
    # 顶部小标
    txt(slide, 0.6, 0.15, 3.0, 0.35,
        f"{page_num:02d} / {total:02d}", 12, DARK_GRAY, False, PP_ALIGN.LEFT, "Segoe UI")
    if section_label:
        txt(slide, 3.0, 0.17, 9.5, 0.35, section_label, 12, TECH_BLUE, False)
    
    bar(slide, 0.6, 0.6, 12.133, 0.008, RGBColor(0x1E, 0x29, 0x3B))
    
    # === 标题 — 大字占足空间 ===
    subtitle = (scene.get('subtitle') or '').strip()
    title = ""
    if subtitle and subtitle != '/':
        first_line = subtitle.split('\n')[0].strip()
        first_line = first_line.replace('【下阶段】', '')
        # 标题控制在20字内，太大就截断
        if len(first_line) > 30:
            first_line = first_line[:30] + "..."
        title = first_line
    
    if not title:
        nar = (scene.get('narration') or '').strip()
        if nar and nar != '/':
            # 取核心句
            parts = nar.split('。')
            title = parts[0].strip()
            if len(title) > 40:
                title = title[:40] + "..."
        else:
            title = f'镜号 {scene["num"]}' 
    
    # 标题 — 大字显示
    txt(slide, 0.6, 0.9, 12.133, 1.2, title, 42, WHITE, True)
    bar(slide, 0.6, 2.2, 3.5, 0.03)
    
    # === 要点卡片 — 大字，撑满剩余空间 ===
    key_points = extract_key_points(subtitle, 4) if subtitle else []
    if not key_points:
        nar = scene['narration'].strip()
        if nar and nar != '/':
            key_points = extract_key_points(nar, 4)
    
    if key_points:
        n = len(key_points)
        card_y = 2.6
        card_total_h = 6.4 - card_y  # 底部留一点边距
        gap = 0.15
        
        if n <= 2:
            # 竖排大卡片
            per_h = (card_total_h - gap * (n - 1)) / n
            for i, pt in enumerate(key_points):
                cy = card_y + i * (per_h + gap)
                card(slide, 0.6, cy, 12.133, per_h)
                txt(slide, 1.0, cy + 0.3, 1.0, 0.6,
                    f"{i+1:02d}", 36, TECH_BLUE, True, PP_ALIGN.LEFT, "Segoe UI")
                # 文字撑满卡片
                txt(slide, 2.0, cy + 0.25, 10.3, per_h - 0.6,
                    pt, 28 if len(pt) < 30 else 24, WHITE, False)
        else:
            # 2x2 网格
            per_w = (12.133 - gap) / 2
            per_h = (card_total_h - gap) / 2
            for i, pt in enumerate(key_points):
                row = i // 2
                col = i % 2
                cx = 0.6 + col * (per_w + gap)
                cy = card_y + row * (per_h + gap)
                card(slide, cx, cy, per_w, per_h)
                txt(slide, cx + 0.3, cy + 0.2, 1.0, 0.6,
                    f"{i+1:02d}", 32, TECH_BLUE, True, PP_ALIGN.LEFT, "Segoe UI")
                txt(slide, cx + 1.5, cy + 0.15, per_w - 2.0, per_h - 0.4,
                    pt, 24 if len(pt) < 30 else 20, WHITE, False)
    
    # 底部旁白简句
    nar = scene['narration'].strip()
    if nar and nar != '/':
        short_nar = nar.split('。')[0] if '。' in nar else nar[:100]
        if len(short_nar) > 100:
            short_nar = short_nar[:100] + "..."
        rect(slide, 0, 6.6, 13.333, 0.9, BG_CARD)
        txt(slide, 0.6, 6.65, 1.2, 0.4, "旁  白", 11, TECH_BLUE, True)
        txt(slide, 1.8, 6.65, 11.0, 0.8, short_nar, 13, GRAY)


def make_ending(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, 13.333, 7.5, BG_DARK)
    bar(slide, 0, 0, 13.333, 0.05)
    
    txt(slide, 0.6, 2.0, 12.1, 1.5, "持续数智创新\n筑牢金融安全屏障", 50, WHITE, True, PP_ALIGN.CENTER)
    bar(slide, 5.0, 3.8, 3.333, 0.03)
    txt(slide, 0.6, 4.2, 12.1, 0.8, "金融让生活更美好", 30, TECH_BLUE, False, PP_ALIGN.CENTER)
    txt(slide, 0.6, 5.3, 12.1, 0.5, "上海银行  ·  安全保卫部", 18, DARK_GRAY, False, PP_ALIGN.CENTER)


def main():
    scenes = parse_scenes()
    print(f"📋 {len(scenes)} 个场景")
    
    prs = Presentation()
    prs.slide_width = Emu(int(SLIDE_W))
    prs.slide_height = Emu(int(SLIDE_H))
    
    total_pages = 20
    p = 0
    
    # 封面
    print("[1] 封面")
    make_cover(prs); p += 1
    
    # === Ch1 ===
    print("[2-6] Ch1 开篇")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    make_chapter(s, 1, "开篇 · 背景与发展概况",
                 "模拟安防 → 数字安防 → 智能应用 → 数智化阶段")
    p += 1
    for sc in scenes:
        if 1 <= sc['num'] <= 4:
            make_content(prs, p, total_pages, sc, "第一章 开篇与背景"); p += 1
    
    # === Ch2 ===
    print("[7-8] Ch2 规划框架")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    make_chapter(s, 2, "整体规划框架", '三大方向 · "看得到 管得牢 用得优"'); p += 1
    for sc in scenes:
        if sc['num'] == 5:
            make_content(prs, p, total_pages, sc, "第二章 规划框架"); p += 1
    
    # === Ch3 ===
    print("[9-18] Ch3 核心成果")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    make_chapter(s, 3, "核心实践成果", '三大板块 · 全面落地'); p += 1
    for sc in scenes:
        if 6 <= sc['num'] <= 13:
            sec = sc.get('section', '') or ''
            label = "第三章 核心成果"
            if '看得到' in sec: label = "板块一 看得到·看得清·看得全"
            elif '防得住' in sec: label = "板块二 防得住·管得牢·抓得实"
            elif '用得优' in sec: label = "板块三 用得优·用得顺"
            make_content(prs, p, total_pages, sc, label); p += 1
    
    # === Ch4 ===
    print("[19-20] Ch4 总结")
    s = prs.slides.add_slide(prs.slide_layouts[6])
    make_chapter(s, 4, "上篇 · 总结展望", "立足安防 · 拥抱AI · 服务全行"); p += 1
    for sc in scenes:
        if sc['num'] == 15:
            make_content(prs, p, total_pages, sc, "第四章 总结展望"); p += 1
    
    # 结尾
    print("[21] 结尾")
    make_ending(prs); p += 1
    
    out = str(OUTPUT_DIR / "上海银行安防宣传片_PPT_v3.pptx")
    prs.save(out)
    print(f"\n✅ {out}  |  {os.path.getsize(out)/1024:.0f} KB  |  {p} 页")


if __name__ == '__main__':
    main()

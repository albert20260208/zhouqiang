#!/usr/bin/env python3
"""镜6 三个要点 PPT —— 精装版 · 纯内容不留空"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import math, os

BLUE   = RGBColor(0x00, 0x2D, 0x73)
BLUE2  = RGBColor(0x00, 0x40, 0x90)
BLUE_L = RGBColor(0xE0, 0xE8, 0xF2)
SILVER = RGBColor(0xA9, 0xAE, 0xB8)
SILVER_L=RGBColor(0xE5, 0xE8, 0xEC)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1C, 0x1C, 0x1C)
GRAY   = RGBColor(0x70, 0x75, 0x7D)
LIGHT  = RGBColor(0xF7, 0xF8, 0xFA)
CARD   = RGBColor(0xFC, 0xFC, 0xFD)

W = Inches(16); H = Inches(9)

prs = Presentation()
prs.slide_width  = W; prs.slide_height = H

def bg(s,c): s.background.fill.solid(); s.background.fill.fore_color.rgb=c

def rc(s,l,t,w,h,f,ln=None,lw=None):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h)
    sh.fill.solid();sh.fill.fore_color.rgb=f
    if ln: sh.line.color.rgb=ln;sh.line.width=Pt(lw or 1)
    else: sh.line.fill.background()
    return sh

def ci(s,l,t,sz,f,ln=None,lw=1):
    sh=s.shapes.add_shape(MSO_SHAPE.OVAL,l,t,sz,sz)
    sh.fill.solid();sh.fill.fore_color.rgb=f
    if ln: sh.line.color.rgb=ln;sh.line.width=Pt(lw)
    else: sh.line.fill.background()
    return sh

def tx(s,l,t,w,h,text,fs=18,c=DARK,b=False,a=PP_ALIGN.LEFT):
    bx=s.shapes.add_textbox(l,t,w,h);tf=bx.text_frame;tf.word_wrap=True
    p=tf.paragraphs[0];p.text=text;p.font.size=Pt(fs);p.font.color.rgb=c
    p.font.bold=b;p.font.name='Microsoft YaHei';p.alignment=a
    return bx

def mtx(s,l,t,w,h,lines,ls=1.2):
    bx=s.shapes.add_textbox(l,t,w,h);tf=bx.text_frame;tf.word_wrap=True
    for i,ln in enumerate(lines):
        txt=ln[0];b=ln[1]if len(ln)>1 else False
        fs=ln[2]if len(ln)>2 else 13;c=ln[3]if len(ln)>3 else DARK
        p=tf.paragraphs[0]if i==0 else tf.add_paragraph()
        p.text=txt;p.font.size=Pt(fs);p.font.color.rgb=c
        p.font.bold=b;p.font.name='Microsoft YaHei';p.space_after=Pt(fs*(ls-1))
    return bx

# ═══════════════════════════════════ P1 ═══════════════════════════════════
s1=prs.slides.add_slide(prs.slide_layouts[6]);bg(s1,WHITE)
# 顶条
rc(s1,Inches(0),Inches(0),W,Inches(0.05),BLUE)
# 01 水印
tx(s1,Inches(12.5),Inches(0.3),Inches(3),Inches(1.8),"01",fs=130,c=BLUE_L,b=True,a=PP_ALIGN.RIGHT)
# 标题
tx(s1,Inches(1.2),Inches(1.5),Inches(10),Inches(1.0),
   "业内率先实现上海地区安防专网全覆盖",fs=38,c=BLUE,b=True)
tx(s1,Inches(1.2),Inches(2.6),Inches(8),Inches(0.5),
   "实现监控中心的\u201c多场所\u201d",fs=22,c=BLUE,b=True)
# 分隔线
rc(s1,Inches(1.2),Inches(3.3),Inches(3.5),Inches(0.02),SILVER)

# 三点列表
pts1=[("▎安防专网","上海地区全覆盖，业内率先完成部署"),
      ("▎多场所监控中心","打破物理空间限制，实现远程集中值守"),
      ("▎核心价值","统一管控 · 降低人力成本 · 提升响应速度")]
for i,(t,d) in enumerate(pts1):
    y=Inches(3.8)+i*Inches(1.2)
    tx(s1,Inches(1.5),y,Inches(5.5),Inches(0.4),t,fs=18,c=BLUE,b=True)
    tx(s1,Inches(1.5),y+Inches(0.4),Inches(5.5),Inches(0.5),d,fs=14,c=GRAY)

# 右侧拓扑图 —— 圆形节点网
cx=Inches(10.8);cy=Inches(5.3)
# 外虚线（小圆点）
for a in range(0,360,20):
    r=Inches(2.2);theta=math.radians(a)
    ci(s1,cx+r*math.cos(theta)-Inches(0.025),cy+r*math.sin(theta)-Inches(0.025),Inches(0.05),BLUE_L)
# 内环
ring=s1.shapes.add_shape(MSO_SHAPE.OVAL,cx-Inches(1.8),cy-Inches(1.8),Inches(3.6),Inches(3.6))
ring.fill.background();ring.line.color.rgb=BLUE_L;ring.line.width=Pt(1.5)
# 中心
ci(s1,cx-Inches(0.55),cy-Inches(0.55),Inches(1.1),BLUE)
tx(s1,cx-Inches(0.55),cy-Inches(0.55),Inches(1.1),Inches(1.1),"监控\n中心",fs=13,c=WHITE,b=True,a=PP_ALIGN.CENTER)
# 6卫星
for i in range(6):
    a=math.radians(-90+i*60);r=Inches(2.2)
    nx=cx+r*math.cos(a)-Inches(0.22);ny=cy+r*math.sin(a)-Inches(0.22)
    ci(s1,nx,ny,Inches(0.44),SILVER_L,BLUE,1.5)
    tx(s1,nx,ny,Inches(0.44),Inches(0.44),f"网点{i+1}",fs=7,c=BLUE,b=True,a=PP_ALIGN.CENTER)

# ═══════════════════════════════════ P2 ═══════════════════════════════════
s2=prs.slides.add_slide(prs.slide_layouts[6]);bg(s2,WHITE)
rc(s2,Inches(0),Inches(0),W,Inches(0.05),BLUE)
tx(s2,Inches(12.5),Inches(0.3),Inches(3),Inches(1.8),"02",fs=130,c=BLUE_L,b=True,a=PP_ALIGN.RIGHT)
tx(s2,Inches(1.2),Inches(1.5),Inches(10),Inches(1.0),
   "业内首家对办公场所、重控场所重点部位\n24 小时监测",fs=36,c=BLUE,b=True)
tx(s2,Inches(1.2),Inches(2.8),Inches(8),Inches(0.5),
   "发现问题实时通知整改  ·  全流程闭环管理",fs=20,c=BLUE,b=True)

# 三列大卡片
cards2=[("全面\n感知","办公场所 · 重控场所 · 燃气阀门\n……物联网传感器实时采集数据",BLUE),
        ("实时\n告警","异常行为自动识别 · 风险事件秒级预警\n多渠道通知到人 · 精准定位风险点",BLUE2),
        ("闭环\n整改","发现→通知→处置→复核→归档\n全流程线上流转 · 可追溯可审计",RGBColor(0x00,0x55,0xA0))]
cw=Inches(3.8);ch=Inches(4.2);g=Inches(0.5);sx=Inches(1.4);cy2=Inches(3.8)
for i,(t,d,ac) in enumerate(cards2):
    x=sx+i*(cw+g)
    rc(s2,x,cy2,cw,ch,CARD,SILVER_L,1)
    # 左侧色条
    rc(s2,x,cy2,Inches(0.08),ch,ac)
    # 标题
    tx(s2,x+Inches(0.35),cy2+Inches(0.8),cw-Inches(0.6),Inches(0.9),t,fs=22,c=ac,b=True)
    # 描述
    tx(s2,x+Inches(0.35),cy2+Inches(2.0),cw-Inches(0.6),Inches(1.8),d,fs=13,c=GRAY)
    # 编号角标
    tx(s2,x+cw-Inches(0.8),cy2+Inches(0.2),Inches(0.6),Inches(0.5),f"0{i+1}",fs=28,c=SILVER_L,b=True,a=PP_ALIGN.RIGHT)

# ═══════════════════════════════════ P3 ═══════════════════════════════════
s3=prs.slides.add_slide(prs.slide_layouts[6]);bg(s3,WHITE)
rc(s3,Inches(0),Inches(0),W,Inches(0.05),BLUE)
tx(s3,Inches(12.5),Inches(0.3),Inches(3),Inches(1.8),"03",fs=130,c=BLUE_L,b=True,a=PP_ALIGN.RIGHT)
tx(s3,Inches(1.2),Inches(1.5),Inches(10),Inches(0.8),
   "搭建业内首个网络安全管理平台",fs=38,c=BLUE,b=True)
tx(s3,Inches(1.2),Inches(2.4),Inches(8),Inches(0.5),
   "六大功能监测模型  ·  全维度安全防护体系",fs=20,c=BLUE,b=True)

# 中心盾牌
cx3=Inches(8);cy3=Inches(5.5)
ci(s3,cx3-Inches(0.9),cy3-Inches(0.9),Inches(1.8),BLUE)
ci(s3,cx3-Inches(0.7),cy3-Inches(0.7),Inches(1.4),WHITE,BLUE,2.5)
tx(s3,cx3-Inches(0.7),cy3-Inches(0.7),Inches(1.4),Inches(1.4),
   "安防\n网络\n安全",fs=16,c=BLUE,b=True,a=PP_ALIGN.CENTER)

# 外环
or_=s3.shapes.add_shape(MSO_SHAPE.OVAL,cx3-Inches(2.4),cy3-Inches(2.4),Inches(4.8),Inches(4.8))
or_.fill.background();or_.line.color.rgb=BLUE_L;or_.line.width=Pt(1)

# 6 功能节点
f6=[("在线准入","设备身份验证\n安全接入控制"),
    ("防外部攻击","防火墙策略\n入侵检测防御"),
    ("隔离内网侵袭","微隔离技术\n横向移动防护"),
    ("扫描系统漏洞","脆弱性自动巡检\n风险评级修复"),
    ("网络稳定性监测","流量行为分析\n异常精准定位"),
    ("入侵预警","实时威胁告警\n自动响应处置")]
for i,(t,d) in enumerate(f6):
    a=math.radians(-90+i*60);r=Inches(3.0)
    nx=cx3+r*math.cos(a)-Inches(0.58);ny=cy3+r*math.sin(a)-Inches(0.58)
    ci(s3,nx,ny,Inches(1.16),CARD,BLUE,1.5)
    rc(s3,nx+Inches(0.12),ny+Inches(0.08),Inches(0.92),Inches(0.035),BLUE)
    tx(s3,nx+Inches(0.08),ny+Inches(0.18),Inches(1.0),Inches(0.35),t,fs=10,c=BLUE,b=True,a=PP_ALIGN.CENTER)
    tx(s3,nx+Inches(0.08),ny+Inches(0.55),Inches(1.0),Inches(0.5),d,fs=8,c=GRAY,a=PP_ALIGN.CENTER)
    # 连线
    sx_=cx3+Inches(0.6)*math.cos(a);sy_=cy3+Inches(0.6)*math.sin(a)
    conn=s3.shapes.add_connector(1,sx_,sy_,nx+Inches(0.58),ny+Inches(0.58))
    conn.line.color.rgb=BLUE_L;conn.line.width=Pt(1.5)

# ── 保存 ──
out='/root/.openclaw/workspace/projects/security-video/output/镜6_三个要点_PPT_v3.pptx'
os.makedirs(os.path.dirname(out),exist_ok=True)
prs.save(out)
print(f'✅ {out}')

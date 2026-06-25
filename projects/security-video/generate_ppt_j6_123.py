#!/usr/bin/env python3
"""镜6 三个要点 PPT —— 上海银行主题色"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── 品牌色 ──
BLUE  = RGBColor(0x00, 0x2D, 0x73)
SILVER = RGBColor(0xA9, 0xAE, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF0, 0xF2, 0xF5)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
ACCENT_LINE = RGBColor(0x00, 0x5A, 0xB3)

prs = Presentation()
prs.slide_width  = Inches(16)
prs.slide_height = Inches(9)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=16, color=DARK_TEXT, line_spacing=1.5):
    """lines: list of (text, bold, font_size_override, color_override)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            text, bold, fs, clr = line, False, font_size, color
        else:
            text = line[0]
            bold = line[1] if len(line) > 1 else False
            fs = line[2] if len(line) > 2 else font_size
            clr = line[3] if len(line) > 3 else color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(fs)
        p.font.color.rgb = clr
        p.font.bold = bold
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(fs * (line_spacing - 1))
    return txBox

def add_number_circle(slide, left, top, size, number, fill_color=BLUE, text_color=WHITE):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(size / Inches(1) * 72 * 0.45)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.CENTER
    return shape

# ═══════════════════════════════════════════
# PAGE 1 — 安防专网全覆盖
# ═══════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide1, WHITE)

# 顶部蓝色装饰条
add_rect(slide1, Inches(0), Inches(0), Inches(16), Inches(0.08), BLUE)

# 左侧蓝色竖条 + 编号
add_rect(slide1, Inches(0.8), Inches(1.2), Inches(0.06), Inches(1.6), BLUE)
add_number_circle(slide1, Inches(0.55), Inches(1.15), Inches(0.55), "01", BLUE, WHITE)

# 标题
add_textbox(slide1, Inches(1.5), Inches(1.15), Inches(11), Inches(0.7),
            "业内率先实现上海地区安防专网全覆盖", 
            font_size=32, color=BLUE, bold=True)

# 副标题
add_textbox(slide1, Inches(1.5), Inches(1.85), Inches(11), Inches(0.5),
            "实现监控中心的多场所", 
            font_size=22, color=SILVER, bold=False)

# 主体 —— 左侧文字 + 右侧图形
# 左侧要点
add_multiline(slide1, Inches(1.5), Inches(2.8), Inches(7), Inches(3.5),
    [
        ("▎ 安防专网", True, 18, BLUE),
        ("上海地区全覆盖，业内率先完成部署", False, 15, DARK_TEXT),
        ("", False, 10, DARK_TEXT),
        ("▎ 多场所监控中心", True, 18, BLUE),
        ("打破物理空间限制，实现远程集中值守", False, 15, DARK_TEXT),
        ("", False, 10, DARK_TEXT),
        ("▎ 核心价值", True, 18, BLUE),
        ("统一管控 · 降低人力成本 · 提升响应速度", False, 15, SILVER),
    ], line_spacing=1.3)

# 右侧 —— 网络拓扑示意（圆形节点连线）
# 中心节点
cx, cy = Inches(11.5), Inches(4.2)
add_number_circle(slide1, cx - Inches(0.55), cy - Inches(0.55), Inches(1.1), "监控\n中心", BLUE, WHITE)
# 周围节点
for angle, label in [(0, "分行1"), (60, "分行2"), (120, "分行3"), (180, "分行4"), (240, "分行5"), (300, "分行6")]:
    import math
    r = Inches(2.0)
    nx = cx + r * math.cos(math.radians(angle)) - Inches(0.35)
    ny = cy + r * math.sin(math.radians(angle)) - Inches(0.35)
    add_number_circle(slide1, nx, ny, Inches(0.7), label, RGBColor(0xE8, 0xED, 0xF2), BLUE)
    # 改文字大小
    for sp in slide1.shapes:
        pass

# 底部下阶段飘字
add_rect(slide1, Inches(0), Inches(8.3), Inches(16), Inches(0.7), LIGHT_BG)
add_textbox(slide1, Inches(1.5), Inches(8.35), Inches(13), Inches(0.5),
            "【下阶段】加快推进各地分行安防专网建设",
            font_size=14, color=SILVER, bold=False, alignment=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════
# PAGE 2 — 重点部位24小时监测
# ═══════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, WHITE)

# 顶部蓝色装饰条
add_rect(slide2, Inches(0), Inches(0), Inches(16), Inches(0.08), BLUE)

# 标题区
add_rect(slide2, Inches(0.8), Inches(1.0), Inches(0.06), Inches(1.4), BLUE)
add_number_circle(slide2, Inches(0.55), Inches(0.95), Inches(0.55), "02", BLUE, WHITE)
add_textbox(slide2, Inches(1.5), Inches(0.95), Inches(11), Inches(0.7),
            "业内首家对办公场所、重控场所重点部位 24 小时监测",
            font_size=30, color=BLUE, bold=True)
add_textbox(slide2, Inches(1.5), Inches(1.6), Inches(11), Inches(0.4),
            "发现问题实时通知整改",
            font_size=20, color=SILVER)

# 三列卡片布局
card_w = Inches(4.0)
card_h = Inches(4.5)
card_y = Inches(2.5)
gap = Inches(0.4)
start_x = Inches(1.2)

cards = [
    ("🔍", "全面感知", "办公场所 · 重控场所\n燃气阀门 · 重点部位\n物联网传感器实时采集"),
    ("⚡", "实时告警", "异常行为自动识别\n风险事件秒级预警\n多渠道通知到人"),
    ("✅", "闭环整改", "问题发现 → 通知 → 处置\n→ 复核 → 归档\n全流程线上追溯"),
]

for i, (icon, title, desc) in enumerate(cards):
    x = start_x + i * (card_w + gap)
    # 卡片背景
    add_rect(slide2, x, card_y, card_w, card_h, LIGHT_BG, ACCENT_LINE)
    # 顶部色条
    add_rect(slide2, x, card_y, card_w, Inches(0.06), BLUE)
    # 图标
    add_textbox(slide2, x + Inches(0.3), card_y + Inches(0.5), card_w - Inches(0.6), Inches(0.8),
                icon, font_size=36, color=BLUE, alignment=PP_ALIGN.LEFT)
    # 标题
    add_textbox(slide2, x + Inches(0.3), card_y + Inches(1.3), card_w - Inches(0.6), Inches(0.5),
                title, font_size=22, color=BLUE, bold=True)
    # 描述
    add_textbox(slide2, x + Inches(0.3), card_y + Inches(2.0), card_w - Inches(0.6), Inches(2.2),
                desc, font_size=14, color=DARK_TEXT)

# 底部
add_rect(slide2, Inches(0), Inches(7.6), Inches(16), Inches(0.06), SILVER)
add_textbox(slide2, Inches(1.2), Inches(7.8), Inches(13.6), Inches(0.5),
            "涵盖：办公场所 | 重控场所 | 燃气阀门 | 重点部位  ·  监测方式：物联网传感器 + AI视觉识别 + 人工复核",
            font_size=12, color=SILVER, alignment=PP_ALIGN.CENTER)

# 底部下阶段
add_rect(slide2, Inches(0), Inches(8.3), Inches(16), Inches(0.7), LIGHT_BG)
add_textbox(slide2, Inches(1.5), Inches(8.35), Inches(13), Inches(0.5),
            "【下阶段】加快推进各地分行安防专网建设",
            font_size=14, color=SILVER)

# ═══════════════════════════════════════════
# PAGE 3 — 网络安全管理平台 6大功能
# ═══════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, WHITE)

# 顶部蓝色装饰条
add_rect(slide3, Inches(0), Inches(0), Inches(16), Inches(0.08), BLUE)

# 标题
add_rect(slide3, Inches(0.8), Inches(1.0), Inches(0.06), Inches(1.4), BLUE)
add_number_circle(slide3, Inches(0.55), Inches(0.95), Inches(0.55), "03", BLUE, WHITE)
add_textbox(slide3, Inches(1.5), Inches(0.95), Inches(11), Inches(0.7),
            "搭建业内首个网络安全管理平台",
            font_size=32, color=BLUE, bold=True)
add_textbox(slide3, Inches(1.5), Inches(1.6), Inches(11), Inches(0.4),
            "六大功能监测模型，全维度防护",
            font_size=20, color=SILVER)

# 中心辐射布局 —— 中间盾牌 + 周围6功能
center_x = Inches(8)
center_y = Inches(5.0)

# 中心 —— 盾牌圆
center_shape = add_number_circle(slide3, center_x - Inches(1.1), center_y - Inches(1.1), Inches(2.2), 
                                   "安防\n网络\n安全", BLUE, WHITE)
# 外环
from pptx.enum.shapes import MSO_SHAPE
ring = slide3.shapes.add_shape(MSO_SHAPE.OVAL, 
    center_x - Inches(1.35), center_y - Inches(1.35), Inches(2.7), Inches(2.7))
ring.fill.background()
ring.line.color.rgb = BLUE
ring.line.width = Pt(2)

# 6个功能节点环形排列
funcs = [
    ("在线准入", "设备接入\n身份验证"),
    ("防外部攻击", "防火墙\n入侵检测"),
    ("隔离内网\n侵袭", "微隔离\n横向防护"),
    ("扫描系统\n漏洞", "脆弱性\n自动巡检"),
    ("监测网络\n稳定性", "流量分析\n异常定位"),
    ("入侵预警", "实时告警\n自动响应"),
]

import math
for i, (title, desc) in enumerate(funcs):
    angle = math.radians(-90 + i * 60)  # start from top
    r = Inches(2.8)
    nx = center_x + r * math.cos(angle) - Inches(0.65)
    ny = center_y + r * math.sin(angle) - Inches(0.65)
    
    # 节点圆
    node = slide3.shapes.add_shape(MSO_SHAPE.OVAL, nx, ny, Inches(1.3), Inches(1.3))
    node.fill.solid()
    node.fill.fore_color.rgb = LIGHT_BG
    node.line.color.rgb = BLUE
    node.line.width = Pt(1.5)
    
    tf = node.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(11)
    p.font.color.rgb = BLUE
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(8)
    p2.font.color.rgb = SILVER
    p2.font.name = 'Microsoft YaHei'
    p2.alignment = PP_ALIGN.CENTER

# 底部
add_rect(slide3, Inches(0), Inches(8.3), Inches(16), Inches(0.7), LIGHT_BG)
add_textbox(slide3, Inches(1.5), Inches(8.35), Inches(13), Inches(0.5),
            "【下阶段】加快推进各地分行安防专网建设",
            font_size=14, color=SILVER)

# ── 保存 ──
out = '/root/.openclaw/workspace/projects/security-video/output/镜6_三个要点_PPT.pptx'
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print(f'✅ 已保存: {out}')

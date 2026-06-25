#!/usr/bin/env python3
"""
上海银行安防宣传片 PPT v4 — 每页独立设计，不套模板
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import openpyxl, os, re
from pathlib import Path

EXCEL_PATH = "/root/.openclaw/media/qqbot/downloads/20260616上海银行安防数字化探索实践短片脚本-V1_1781602866070_644fd5.xlsx"
OUTPUT_DIR = Path("/root/.openclaw/workspace/projects/security-video/output")
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# 配色
BG  = RGBColor(0x08, 0x0C, 0x14)
BG2 = RGBColor(0x10, 0x15, 0x20)
BL  = RGBColor(0x3B, 0x82, 0xF6)  # 主蓝
BL2 = RGBColor(0x60, 0xA5, 0xFA)  # 浅蓝
CY  = RGBColor(0x22, 0xD3, 0xEE)  # 青
W   = RGBColor(0xFF, 0xFF, 0xFF)
G   = RGBColor(0x9C, 0xA3, 0xAF)
DG  = RGBColor(0x6B, 0x72, 0x80)

SW = Inches(13.333); SH = Inches(7.5)

def r(slide, l, t, w, h, c, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()
    return s

def card(slide, l, t, w, h, c=BG2): return r(slide, l, t, w, h, c, MSO_SHAPE.ROUNDED_RECTANGLE)

def t(slide, l, t, w, h, text, size=22, color=W, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color; p.font.name = "微软雅黑"
    p.alignment = align
    return box

def tl(slide, l, t, w, h, lines, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    for i, (text, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.font.size = Pt(size); p.font.bold = bold
        p.font.color.rgb = color; p.font.name = "微软雅黑"; p.alignment = align
        p.space_after = Pt(4)
    return box

def bar(slide, l, t, w, h=0.02, c=BL):
    r(slide, l, t, w, h, c)

def tbar(slide): bar(slide, 0, 0, 13.333, 0.04)

def page_num(slide, n, total):
    t(slide, 0.5, 0.15, 2, 0.3, f"{n:02d} / {total:02d}", 10, DG, False, PP_ALIGN.LEFT)

# ==================== 场景解析 ====================
def parse():
    wb = openpyxl.load_workbook(EXCEL_PATH)
    ws = wb['上篇']
    scenes, chap, sec = [], "", ""
    for row in ws.iter_rows(values_only=True):
        vals = [str(c).strip() if c else '' for c in row[:6]]
        f = vals[0]
        if all(v == '' for v in vals[:4]): continue
        if '章节' in f: chap = f; continue
        if '板块' in f: sec = f; continue
        if '部分' in f or '镜号' in f or '此篇' in f: continue
        if f.isdigit():
            scenes.append(dict(num=int(f), time=vals[1],
                visual=vals[2], subtitle=vals[3], narration=vals[4],
                materials=vals[5] if len(vals)>5 else '',
                chapter=chap, section=sec))
    return scenes

def pts(text, n=5):
    """提取要点"""
    if not text or text == '/': return []
    out, seen = [], set()
    for line in text.replace('<br>','\n').split('\n'):
        line = line.strip()
        if not line or '下阶段' in line: continue
        m = re.match(r'^\d+[、.]?\s*(.+)', line)
        if m: line = m.group(1)
        if len(line) > 3 and line not in seen:
            seen.add(line); out.append(line)
    return out[:n]


# ==================== 页面设计 ====================

def cover(prs):
    """封面"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r(s,0,0,13.333,7.5,BG); tbar(s)
    # 左侧装饰块
    r(s,0,0,0.25,7.5,BL)
    t(s,1.5,1.8,10.5,1.5,"上海银行",60,W,True)
    t(s,1.5,3.0,10.5,1.0,"安防数智化应用展示",44,BL2,True)
    bar(s,1.5,4.2,3.5)
    t(s,1.5,4.6,10.5,0.7,"统筹发展与安全，筑牢金融安全防线",24,G)
    t(s,1.5,5.5,10.5,0.4,"上篇 · 安防数智化探索实践",16,DG)


def scene1_title(prs, pn, total, sc):
    """镜号1: LOGO + 标题 — 居中简洁"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r(s,0,0,13.333,7.5,BG); tbar(s); page_num(s,pn,total)
    # 大 LOGO 区域
    r(s,5.0,1.5,3.333,1.8,BG2)
    t(s,5.6,2.0,2.133,0.8,"SHANGHAI\nBANK",18,W,True,PP_ALIGN.CENTER)
    t(s,5.0,3.5,3.333,0.5,"上海银行",28,BL,True,PP_ALIGN.CENTER)
    bar(s,4.5,4.3,4.333,0.02)
    t(s,2.5,4.8,8.333,1.0,"安防数智化应用展示",42,W,True,PP_ALIGN.CENTER)
    t(s,2.0,5.8,9.333,0.6,"统筹发展与安全，筑牢金融安全防线",20,G,False,PP_ALIGN.CENTER)


def scene2_intro(prs, pn, total, sc):
    """镜号2: 航拍+简介 — 左图右文"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r(s,0,0,13.333,7.5,BG); tbar(s); page_num(s,pn,total)
    
    # 左: 占位图区域
    r(s,0.5,0.8,5.5,5.8,BG2)
    t(s,1.5,2.5,3.5,1.0,"素材1-1\n总行浦西大楼航拍",16,DG,False,PP_ALIGN.CENTER)
    t(s,1.5,3.5,3.5,1.0,"素材1-2\n金库视频",16,DG,False,PP_ALIGN.CENTER)
    
    # 右: 文字
    t(s,6.5,1.0,6.5,1.2,"背景与发展概况",36,W,True)
    bar(s,6.5,2.3,2.5)
    pts_text = pts(sc['subtitle'])
    if pts_text:
        tl(s,6.5,2.7,6.5,4.0,
           [(p,20,W,False) for p in pts_text])


def scene3_timeline(prs, pn, total, sc):
    """镜号3: 四大发展阶段 — 横向时间轴"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r(s,0,0,13.333,7.5,BG); tbar(s); page_num(s,pn,total)
    
    t(s,0.5,0.8,12.5,0.8,"安防四大发展阶段",36,W,True)
    bar(s,0.5,1.7,2.5)
    
    stages = [
        ("模拟安防", "GA38-2004\n24h人工轮巡"),
        ("数字安防", "2018-2021\n模改数"),
        ("智能应用", "2021-2025\n智能化功能"),
        ("数智化阶段", "2025→\n大模型驱动"),
    ]
    for i, (title, desc) in enumerate(stages):
        cx = 0.5 + i * 3.1
        color = BL if i == 3 else (BL2 if i == 2 else G)
        card(s, cx, 2.3, 2.9, 3.5)
        t(s, cx+0.3, 2.5, 2.3, 0.6, f"0{i+1}", 32, BL, True, PP_ALIGN.LEFT)
        t(s, cx+0.3, 3.2, 2.3, 0.6, title, 22, W, True)
        t(s, cx+0.3, 3.9, 2.3, 1.2, desc, 14, G)
        # 箭头
        if i < 3:
            t(s, cx+2.9, 3.7, 0.5, 0.5, "→", 28, BL2, True, PP_ALIGN.CENTER)
    
    t(s,0.5,6.2,12.5,0.8,sc['narration'][:150],13,G)


def scene4_techbase(prs, pn, total, sc):
    """镜号4: 三大技术底座 — 三列图标式"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r(s,0,0,13.333,7.5,BG); tbar(s); page_num(s,pn,total)
    
    t(s,0.5,0.8,12.5,0.8,"技术底座 · 三大引擎",36,W,True)
    bar(s,0.5,1.7,2.5)
    
    bases = [
        ("物联网技术", "摄像机、传感器、读识设备\n触探设备等物联部署\n清晰感知、识别、监测", "📡"),
        ("自动化系统", "视频分析、AI识别算法\n风险隐患提前预警\n设备24h不间断监测", "⚙️"),
        ("数据驱动", "安全管理数据结构化\n分析人员行为\n辅助管理决策", "📊"),
    ]
    for i, (title, desc, icon) in enumerate(bases):
        cx = 0.5 + i * 4.2
        card(s, cx, 2.2, 4.0, 4.2)
        t(s, cx+0.3, 2.4, 3.4, 0.5, icon, 36, W, False, PP_ALIGN.LEFT)
        t(s, cx+0.3, 3.1, 3.4, 0.5, title, 24, BL, True)
        bar(s, cx+0.3, 3.7, 1.5, 0.015, BL2)
        tl(s, cx+0.3, 4.0, 3.4, 2.0,
           [(l.strip(),16,G,False) for l in desc.split('\n')])


def scene5_framework(prs, pn, total, sc):
    """镜号5: 看管用框架 — 中心辐射"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r(s,0,0,13.333,7.5,BG); tbar(s); page_num(s,pn,total)
    
    t(s,0.5,0.8,12.5,0.8,"一屏观全行 · 一网管全行",36,W,True)
    bar(s,0.5,1.7,2.5)
    
    # 三个大卡片
    parts = [
        ("看", "3个突破", ["专网建设取得突破", "系统整合取得突破", "数智化监测取得突破"]),
        ("管", "3个创新", ["数智化全流程闭环管理", "全员安防线上履职APP", "安全生产数字画像"]),
        ("用", "2个探索", ["AI模型业务横向经营赋能", "AI模型基层纵向履职赋能"]),
    ]
    for i, (title, sub, items) in enumerate(parts):
        cx = 0.5 + i * 4.2
        card(s, cx, 2.2, 4.0, 4.5)
        t(s, cx+0.4, 2.4, 3.2, 0.8, title, 48, BL, True)
        t(s, cx+0.4, 3.4, 3.2, 0.4, sub, 16, BL2, True)
        bar(s, cx+0.4, 3.9, 1.5, 0.015, BL2)
        for j, item in enumerate(items):
            t(s, cx+0.4, 4.2+j*0.5, 3.2, 0.4, f"▸ {item}", 14, G)


def scene6_8_cards(prs, pn, total, sc, section_label):
    """板块一通用: 标题 + 3卡片 + 数字亮点"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r(s,0,0,13.333,7.5,BG); tbar(s); page_num(s,pn,total)
    
    t(s,0.3,0.15,8,0.3,section_label,11,BL2)
    
    subtitle = sc['subtitle']
    lines = subtitle.replace('<br>','\n').split('\n') if subtitle and subtitle!='/' else []
    main_title = lines[0].strip() if lines else ""
    main_title = main_title.replace('【下阶段】','')
    
    t(s,0.5,0.7,12.5,0.9,main_title,34,W,True)
    bar(s,0.5,1.7,2.5)
    
    p_list = pts(subtitle, 3)
    for i, pt in enumerate(p_list):
        cy = 2.2 + i * 1.5
        card(s, 0.5, cy, 12.333, 1.3)
        t(s, 0.8, cy+0.15, 0.8, 0.6, f"0{i+1}", 28, BL, True, PP_ALIGN.LEFT)
        t(s, 1.8, cy+0.2, 10.5, 0.6, pt, 20, W)
        if i < len(p_list)-1:
            t(s, 12.5, cy+0.45, 0.5, 0.3, "", 12, G, False, PP_ALIGN.RIGHT)
    
    # 底部旁白
    nar = sc['narration']
    if nar and nar != '/':
        short = nar.split('。')[0].strip()
        if len(short) > 120: short = short[:120]+"..."
        r(s,0,6.2,13.333,1.3,BG2)
        t(s,0.5,6.3,1.0,0.3,"旁白",10,BL,True)
        t(s,1.5,6.3,11.3,0.6,short,13,G)


def scene9_11_platform(prs, pn, total, sc, section_label):
    """板块二通用: 左图右文 — 平台展示型"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r(s,0,0,13.333,7.5,BG); tbar(s); page_num(s,pn,total)
    t(s,0.3,0.15,8,0.3,section_label,11,BL2)
    
    subtitle = sc['subtitle']
    lines = subtitle.replace('<br>','\n').split('\n') if subtitle and subtitle!='/' else []
    main_title = lines[0].strip() if lines else ""
    main_title = main_title.replace('【下阶段】','')
    
    t(s,0.5,0.7,12.5,0.9,main_title,34,W,True)
    bar(s,0.5,1.7,2.5)
    
    # 左: 占位图
    r(s,0.5,2.2,5.5,4.5,BG2)
    mat = sc.get('materials','')
    mat_short = mat.split('【')[1].split('】')[0] if '【' in mat else '平台界面截图'
    t(s,1.5,3.5,3.5,1.0,mat_short,14,DG,False,PP_ALIGN.CENTER)
    
    # 右: 要点
    p_list = pts(subtitle, 4)
    tl(s,6.5,2.2,6.5,4.5,
       [(f"▸ {p}",20,W,False) for p in p_list])


def scene12_13_ai(prs, pn, total, sc, section_label):
    """板块三通用: 大字数据 + 影响数字"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r(s,0,0,13.333,7.5,BG); tbar(s); page_num(s,pn,total)
    t(s,0.3,0.15,8,0.3,section_label,11,BL2)
    
    subtitle = sc['subtitle']
    lines = subtitle.replace('<br>','\n').split('\n') if subtitle and subtitle!='/' else []
    main_title = lines[0].strip() if lines else ""
    main_title = main_title.replace('【下阶段】','')
    
    t(s,0.5,0.7,12.5,0.9,main_title,34,W,True)
    bar(s,0.5,1.7,2.5)
    
    # 提取数字
    narration = sc['narration']
    numbers = re.findall(r'(\d{3,4})\s*小时', narration)
    
    p_list = pts(subtitle, 2) or pts(narration, 2)
    
    if numbers:
        # 大字数据展示
        for i, num in enumerate(numbers[:2]):
            cx = 0.5 + i * 6.3
            card(s, cx, 2.2, 6.0, 3.0)
            t(s, cx+0.5, 2.5, 5.0, 1.0, f"{num}小时", 48, BL, True, PP_ALIGN.CENTER)
            if i < len(p_list):
                t(s, cx+0.5, 3.5, 5.0, 0.8, p_list[i], 16, G, False, PP_ALIGN.CENTER)
    else:
        for i, pt in enumerate(p_list):
            cy = 2.2 + i * 2.0
            card(s, 0.5, cy, 12.333, 1.8)
            t(s, 0.8, cy+0.3, 1.0, 0.6, f"0{i+1}", 32, BL, True)
            t(s, 2.0, cy+0.35, 10.5, 1.0, pt, 22, W)


def scene15_closing(prs, pn, total, sc):
    """镜号15: 总结 — 居中简洁"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r(s,0,0,13.333,7.5,BG); tbar(s); page_num(s,pn,total)
    
    t(s,1.0,1.5,11.5,1.2,"安防数智化道阻且长",44,W,True)
    t(s,1.0,2.8,11.5,1.2,"我们将继续秉承",28,G,False)
    t(s,1.0,3.5,11.5,1.2,"立足安防 · 服务全行",36,BL,True)
    bar(s,1.0,4.8,3.0)
    
    ks = ["坚持专业主义、长期主义","拥抱AI，深挖安防数据价值","科技兴安、智慧赋能"]
    for i, k in enumerate(ks):
        t(s,1.0,5.2+i*0.5,11.5,0.4,f"▸ {k}",18,G)
    
    t(s,1.0,6.7,11.5,0.4,"为我行高质量发展保驾护航",20,W,True)


def chapter(prs, num, title, sub=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r(s,0,0,13.333,7.5,BG); tbar(s)
    t(s,0.5,0.8,4,0.5,f"CHAPTER {num:02d}",18,BL,True)
    r(s,0,6.2,13.333,1.3,BG2)
    t(s,0.5,2.2,12.5,1.5,title,48,W,True)
    if sub:
        bar(s,0.5,3.8,3.0)
        t(s,0.5,4.2,12.5,0.8,sub,22,G)


def ending(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r(s,0,0,13.333,7.5,BG); tbar(s)
    r(s,0,0,0.25,7.5,BL)
    t(s,1.5,2.0,10.5,1.2,"持续数智创新\n筑牢金融安全屏障",48,W,True)
    bar(s,1.5,3.6,3.0)
    t(s,1.5,4.0,10.5,0.7,"金融让生活更美好",26,BL2)
    t(s,1.5,5.0,10.5,0.5,"上海银行  ·  安全保卫部",16,DG)


# ==================== 主流程 ====================
def main():
    scenes = parse()
    print(f"📋 {len(scenes)} 个场景")
    
    prs = Presentation()
    prs.slide_width = Emu(int(SW)); prs.slide_height = Emu(int(SH))
    
    total = 21; p = 0
    
    print("封面"); cover(prs); p += 1
    
    # Ch1
    print("Ch1"); chapter(prs,1,"开篇 · 背景与发展概况","模拟→数字→智能→数智化"); p += 1
    
    for sc in scenes:
        if sc['num'] == 1:
            scene1_title(prs,p,total,sc); p += 1
        elif sc['num'] == 2:
            scene2_intro(prs,p,total,sc); p += 1
        elif sc['num'] == 3:
            scene3_timeline(prs,p,total,sc); p += 1
        elif sc['num'] == 4:
            scene4_techbase(prs,p,total,sc); p += 1
    
    # Ch2
    print("Ch2"); chapter(prs,2,"整体规划框架",'三大方向 · "看得到 管得牢 用得优"'); p += 1
    
    for sc in scenes:
        if sc['num'] == 5:
            scene5_framework(prs,p,total,sc); p += 1
    
    # Ch3
    print("Ch3"); chapter(prs,3,"核心实践成果","三大板块 · 全面落地"); p += 1
    
    sec_map = {
        '看得到': '板块一 看得到·看得清·看得全',
        '防得住': '板块二 防得住·管得牢·抓得实',
        '用得优': '板块三 用得优·用得顺',
    }
    for sc in scenes:
        if 6 <= sc['num'] <= 13:
            sec = sc.get('section','') or ''
            label = "核心实践成果"
            for k,v in sec_map.items():
                if k in sec: label = v; break
            
            if sc['num'] >= 12:  # 板块三
                scene12_13_ai(prs,p,total,sc,label)
            elif sc['num'] >= 9:  # 板块二
                scene9_11_platform(prs,p,total,sc,label)
            else:  # 板块一
                scene6_8_cards(prs,p,total,sc,label)
            p += 1
    
    # Ch4
    print("Ch4"); chapter(prs,4,"上篇 · 总结展望","立足安防 · 拥抱AI · 服务全行"); p += 1
    
    for sc in scenes:
        if sc['num'] == 15:
            scene15_closing(prs,p,total,sc); p += 1
    
    print("结尾"); ending(prs); p += 1
    
    out = str(OUTPUT_DIR / "上海银行安防宣传片_PPT_v4.pptx")
    prs.save(out)
    print(f"\n✅ {out}  |  {os.path.getsize(out)/1024:.0f} KB  |  {p} 页")


if __name__ == '__main__':
    main()

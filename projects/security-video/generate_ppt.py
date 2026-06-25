#!/usr/bin/env python3
"""
上海银行安防宣传片 PPT 生成脚本
读取Excel脚本 → 生成专业商务PPTX
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import openpyxl
import os
from pathlib import Path

# ====== 配置 ======
EXCEL_PATH = "/root/.openclaw/media/qqbot/downloads/20260616上海银行安防数字化探索实践短片脚本-V1_1781602866070_644fd5.xlsx"
OUTPUT_DIR = Path("/root/.openclaw/workspace/projects/security-video/output")
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# 配色
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BANK_BLUE = RGBColor(0x00, 0x3D, 0x7C)
TECH_BLUE = RGBColor(0x25, 0x63, 0xEB)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
MID_TEXT = RGBColor(0x4A, 0x4A, 0x4A)
LIGHT_TEXT = RGBColor(0x8A, 0x8A, 0x8A)
ACCENT_GOLD = RGBColor(0xC8, 0x96, 0x2E)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF8)
DIVIDER_COLOR = RGBColor(0xE0, 0xE0, 0xE5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_CN = "微软雅黑"  # Chinese font, fallback
FONT_EN = "Segoe UI"


def parse_scenes():
    """解析Excel脚本"""
    wb = openpyxl.load_workbook(EXCEL_PATH)
    ws = wb['上篇']
    
    scenes = []
    current_chapter = ""
    current_section = ""
    
    for row in ws.iter_rows(values_only=True):
        vals = [str(c).strip() if c else '' for c in row[:6]]
        
        if all(v == '' for v in vals[:4]):
            continue
        
        first = vals[0]
        
        if '章节' in first:
            current_chapter = first
            continue
        if '板块' in first:
            current_section = first
            continue
        if '部分' in first or '镜号' in first or '此篇' in first:
            continue
        
        if first.isdigit():
            time_str = vals[1]
            scenes.append({
                'num': int(first),
                'time': time_str,
                'visual': vals[2].replace('<br>', '\n'),
                'subtitle': vals[3].replace('<br>', '\n'),
                'narration': vals[4].replace('<br>', '\n'),
                'materials': vals[5] if len(vals) > 5 else '',
                'chapter': current_chapter,
                'section': current_section,
            })
    
    return scenes


def add_bg_rect(slide, left, top, width, height, color, alpha=None):
    """添加纯色矩形背景"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        shape.fill.fore_color.brightness = alpha
    return shape


def set_font(run, size_pt=None, bold=None, color=None, font_name=None):
    """设置字体属性"""
    if size_pt:
        run.font.size = Pt(size_pt)
    if bold is not None:
        run.font.bold = bold
    if color:
        run.font.color.rgb = color
    if font_name:
        run.font.name = font_name


def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT_CN):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, font_size=16, color=DARK_TEXT, line_spacing=1.5, alignment=PP_ALIGN.LEFT):
    """添加多行文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, line_data in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        if isinstance(line_data, tuple):
            text, size, bold, clr = line_data
        else:
            text, size, bold, clr = line_data, font_size, False, color
        
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = clr
        p.font.name = FONT_CN
        p.alignment = alignment
        p.space_after = Pt(size * (line_spacing - 1) * 0.5)
    
    return txBox


def create_cover_slide(prs):
    """封面页：上海银行 LOGO + 标题"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, WHITE)
    
    # 顶部蓝色装饰条
    add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BANK_BLUE)
    
    # 左侧蓝色竖条
    add_bg_rect(slide, Inches(0.8), Inches(2.0), Inches(0.06), Inches(3.5), TECH_BLUE)
    
    # 标题
    add_text_box(slide, Inches(1.3), Inches(2.2), Inches(10), Inches(1.2),
                 "上海银行安防数智化应用展示", 44, BANK_BLUE, True)
    
    # 副标题
    add_text_box(slide, Inches(1.3), Inches(3.5), Inches(10), Inches(0.8),
                 "统筹发展与安全，筑牢金融安全防线", 24, MID_TEXT, False)
    
    # 底部装饰线
    add_bg_rect(slide, Inches(1.3), Inches(4.6), Inches(3), Inches(0.03), TECH_BLUE)
    
    # 底部副标题
    add_text_box(slide, Inches(1.3), Inches(5.0), Inches(10), Inches(0.6),
                 "上篇 · 安防数智化探索实践", 16, LIGHT_TEXT, False)


def create_chapter_divider(prs, title, subtitle=""):
    """章节分隔页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BANK_BLUE)
    
    # 淡色叠加
    add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, RGBColor(0x00, 0x50, 0xA0))
    
    add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10.5), Inches(1.2),
                 title, 40, WHITE, True)
    
    if subtitle:
        add_bg_rect(slide, Inches(1.5), Inches(4.2), Inches(2), Inches(0.03), WHITE)
        add_text_box(slide, Inches(1.5), Inches(4.5), Inches(10.5), Inches(0.8),
                     subtitle, 20, RGBColor(0xCC, 0xDD, 0xF0), False)


def create_content_slide(prs, scene, page_num, total_pages, chapter_title=""):
    """内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, WHITE)
    
    # === 顶部区域 ===
    # 蓝色顶部线条
    add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.04), BANK_BLUE)
    
    # 页码+章节
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.4),
                 f"0{page_num}" if page_num < 10 else str(page_num),
                 28, BANK_BLUE, True, PP_ALIGN.LEFT, FONT_EN)
    
    if chapter_title:
        add_text_box(slide, Inches(2.0), Inches(0.35), Inches(8), Inches(0.4),
                     chapter_title, 14, LIGHT_TEXT, False)
    
    # 分隔线
    add_bg_rect(slide, Inches(0.8), Inches(0.85), Inches(11.733), Inches(0.01), DIVIDER_COLOR)
    
    # === 主标题 ===
    subtitle = scene['subtitle'].strip()
    if subtitle and subtitle != '/':
        # 提取核心标题（第一行或最关键的）
        lines = subtitle.split('\n')
        main_title = lines[0].strip()
        # 清理标记
        main_title = main_title.replace('【下阶段】', '')
        if len(main_title) > 60:
            main_title = main_title[:60]
        
        # 标题区域背景色块
        add_bg_rect(slide, Inches(0.8), Inches(1.1), Inches(11.733), Inches(0.0), LIGHT_BG)
        
        add_text_box(slide, Inches(1.0), Inches(1.15), Inches(11.333), Inches(0.7),
                     main_title, 28, BANK_BLUE, True)
    else:
        add_text_box(slide, Inches(1.0), Inches(1.15), Inches(11.333), Inches(0.7),
                     f"镜号 {scene['num']} | {scene['time']}", 24, BANK_BLUE, True)
    
    # === 核心内容区 ===
    # 处理字幕内容（详细描述）
    content_y = Inches(2.1)
    
    if subtitle and subtitle != '/':
        subtitle_text = subtitle.strip()
        
        # 清理
        subtitle_text = subtitle_text.replace('【下阶段】', '\n▸ 下阶段：')
        
        # 拆成要点
        points = []
        for line in subtitle_text.split('\n'):
            line = line.strip()
            if not line:
                continue
            # 数字开头的条目
            if line and (line[0].isdigit() and ('、' in line[:3] or '.' in line[:3])):
                points.append(('  ' + line, 15, False, MID_TEXT))
            elif line.startswith('▸'):
                points.append((line, 14, False, TECH_BLUE))
            else:
                points.append((line, 16, False, DARK_TEXT))
        
        # 限制行数，超出的截断
        if len(points) > 12:
            points = points[:11] + [('  ...', 14, False, LIGHT_TEXT)]
        
        add_multiline_text(slide, Inches(1.0), content_y,
                          Inches(11.333), Inches(4.5),
                          points, line_spacing=1.6)
    
    # === 旁白区域（底部） ===
    narration = scene['narration'].strip()
    if narration and narration != '/':
        # 底部区域
        box_top = Inches(6.3)
        add_bg_rect(slide, Inches(0.8), box_top, Inches(11.733), Inches(0.9), LIGHT_BG)
        
        # 旁白标签
        add_text_box(slide, Inches(1.0), Inches(6.35), Inches(1.5), Inches(0.3),
                     "旁白", 10, TECH_BLUE, True)
        
        # 旁白文本 - 限制长度
        if len(narration) > 200:
            narration = narration[:200] + "..."
        
        add_text_box(slide, Inches(1.0), Inches(6.55), Inches(11.333), Inches(0.6),
                     narration, 11, MID_TEXT, False)
    
    # === 素材提示（右下） ===
    materials = scene.get('materials', '').strip()
    if materials and materials != '/':
        mat_label = "素材"
        add_text_box(slide, Inches(9.5), Inches(0.35), Inches(3.5), Inches(0.4),
                     f"📎 {mat_label}", 10, LIGHT_TEXT, False, PP_ALIGN.RIGHT)
    
    # === 底部页码 ===
    add_text_box(slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.3),
                 f"{page_num} / {total_pages}", 9, LIGHT_TEXT, False, PP_ALIGN.RIGHT, FONT_EN)


def create_ending_slide(prs):
    """结尾页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, WHITE)
    add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BANK_BLUE)
    
    # 中间内容
    add_text_box(slide, Inches(2.5), Inches(2.5), Inches(8.5), Inches(1.0),
                 "持续数智创新  筑牢金融安全屏障", 36, BANK_BLUE, True, PP_ALIGN.CENTER)
    
    add_bg_rect(slide, Inches(5.5), Inches(3.5), Inches(2.333), Inches(0.03), TECH_BLUE)
    
    add_text_box(slide, Inches(2.5), Inches(3.8), Inches(8.5), Inches(0.6),
                 "金融让生活更美好", 22, MID_TEXT, False, PP_ALIGN.CENTER)
    
    add_text_box(slide, Inches(2.5), Inches(4.6), Inches(8.5), Inches(0.5),
                 "上海银行", 20, LIGHT_TEXT, False, PP_ALIGN.CENTER)


def main():
    print("📋 解析Excel脚本...")
    scenes = parse_scenes()
    print(f"   共 {len(scenes)} 个场景\n")
    
    # 分组场景
    groups = {
        'cover': [],
        'ch1': [],  # 开篇 镜号1-4
        'ch2': [],  # 框架 镜号5
        'ch3': [],  # 核心成果 镜号6-13
        'ch4': [],  # 总结 镜号15
    }
    
    for s in scenes:
        n = s['num']
        if n <= 4:
            groups['ch1'].append(s)
        elif n == 5:
            groups['ch2'].append(s)
        elif 6 <= n <= 13:
            groups['ch3'].append(s)
        elif n == 15:
            groups['ch4'].append(s)
    
    # 创建PPT
    prs = Presentation()
    prs.slide_width = Emu(int(SLIDE_W))
    prs.slide_height = Emu(int(SLIDE_H))
    
    page_num = 0
    
    # 封面
    print("🎨 封面页")
    create_cover_slide(prs)
    page_num += 1
    
    # 第一章
    print("🎨 第一章：开篇・背景与发展概况")
    create_chapter_divider(prs, "开篇 · 背景与发展概况", "模拟安防 → 数字安防 → 智能应用 → 数智化阶段")
    page_num += 1
    
    for s in groups['ch1']:
        print(f"   镜号{s['num']} | {s['time']}")
        create_content_slide(prs, s, page_num, len(scenes) + 5, "第一章 开篇与背景")
        page_num += 1
    
    # 第二章
    print("🎨 第二章：整体规划框架")
    create_chapter_divider(prs, "整体规划框架", "看得到 · 管得牢 · 用得优")
    page_num += 1
    
    for s in groups['ch2']:
        print(f"   镜号{s['num']} | {s['time']}")
        create_content_slide(prs, s, page_num, len(scenes) + 5, "第二章 规划框架")
        page_num += 1
    
    # 第三章（带子板块）
    print("🎨 第三章：核心实践成果")
    create_chapter_divider(prs, "核心实践成果", '三大板块 · "看管用"全面落地')
    page_num += 1
    
    for s in groups['ch3']:
        print(f"   镜号{s['num']} | {s['time']}")
        section_label = ""
        if '看得到' in (s.get('section', '') or ''):
            section_label = "板块一 视觉与监测突破"
        elif '防得住' in (s.get('section', '') or ''):
            section_label = "板块二 全流程管理创新"
        elif '用得优' in (s.get('section', '') or ''):
            section_label = "板块三 AI赋能业务与基层"
        
        create_content_slide(prs, s, page_num, len(scenes) + 5, section_label or "第三章 核心成果")
        page_num += 1
    
    # 第四章
    print("🎨 第四章：总结展望")
    create_chapter_divider(prs, "上篇 · 总结展望", "立足安防 · 服务全行")
    page_num += 1
    
    for s in groups['ch4']:
        print(f"   镜号{s['num']} | {s['time']}")
        create_content_slide(prs, s, page_num, len(scenes) + 5, "第四章 总结展望")
        page_num += 1
    
    # 结尾
    print("🎨 结尾页")
    create_ending_slide(prs)
    
    # 保存
    output_path = str(OUTPUT_DIR / "上海银行安防宣传片_PPT.pptx")
    prs.save(output_path)
    
    file_size = os.path.getsize(output_path) / 1024
    print(f"\n✅ 完成! PPT: {output_path}")
    print(f"   大小: {file_size:.0f} KB")
    print(f"   页数: {page_num + 1} 页")
    
    return output_path


if __name__ == '__main__':
    main()

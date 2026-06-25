#!/usr/bin/env python3
"""镜6 三个要点 PPT —— 上海银行主题色 · 精装版"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import math, os

BLUE   = RGBColor(0x00, 0x2D, 0x73)
BLUE2  = RGBColor(0x00, 0x3D, 0x8C)
BLUE_L = RGBColor(0xE8, 0xEF, 0xF5)
SILVER = RGBColor(0xA9, 0xAE, 0xB8)
SILVER_L=RGBColor(0xE0, 0xE3, 0xE8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1A, 0x1A, 0x1A)
GRAY   = RGBColor(0x6B, 0x70, 0x78)
LIGHT  = RGBColor(0xF5, 0xF6, 0xF8)
CARD_BG= RGBColor(0xFA, 0xFB, 0xFC)

prs = Presentation()
prs.slide_width  = Inches(16)
prs.slide_height = Inches(9)

def bg(slide, color): 
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill, line=None, lw=None, alpha=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = Pt(lw or 0.5)
    else: s.line.fill.background()
    if alpha is not None: 
        s.fill.fore_color._color.set(qn('a:alpha'), int(alpha * 1000))
    return s

def circle(slide, l, t, sz, fill, line=None, lw=1):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, sz, sz)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = Pt(lw)
    else: s.line.fill.background()
    return s

def tb(slide, l, t, w, h, text, fs=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, font='Microsoft YaHei'):
    bx = slide.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(fs)
    p.font.color.rgb = color; p.font.bold = bold; p.font.name = font; p.alignment = align
    return bx

def ml(slide, l, t, w, h, lines, ls=1.3):
    bx = slide.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        txt, b, fs, clr = (ln[0], ln[1] if len(ln)>1 else False, ln[2] if len(ln)>2 else 14, ln[3] if len(ln)>3 else DARK) if not isinstance(ln, str) else (ln, False, 14, DARK)
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text = txt; p.font.size = Pt(fs); p.font.color.rgb = clr
        p.font.bold = b; p.font.name = 'Microsoft YaHei'; p.space_after = Pt(fs*(ls-1))
    return bx

def line(slide, x1,y1,x2,y2, color=SILVER, w=1):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # MSO_CONNECTOR.STRAIGHT
    connector.line.color.rgb = color; connector.line.width = Pt(w)

def corner_decor(slide, x, y, size, color=BLUE):
    """L-shaped corner accent"""
    rect(slide, x, y, size, Inches(0.03), color)
    rect(slide, x, y, Inches(0.03), size, color)

# ═══════════════════════════════════════════
# PAGE 1 — 安防专网全覆盖
# ═══════════════════════════════════════════
s1 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s1, WHITE)

# 顶部渐变条（用三层矩形模拟深浅）
rect(s1, Inches(0), Inches(0), Inches(16), Inches(0.06), BLUE)
rect(s1, Inches(0), Inches(0.06), Inches(16), Inches(0.02), BLUE2)

# 左上角装饰
corner_decor(s1, Inches(0.6), Inches(0.5), Inches(0.4), BLUE)
corner_decor(s1, Inches(15.0), Inches(8.1), Inches(0.4), SILVER)

# 超大淡色编号水印
tb(s1, Inches(11.5), Inches(0.3), Inches(4), Inches(2.5), "01", fs=120, color=BLUE_L, bold=True, align=PP_ALIGN.RIGHT)

# 左侧主内容区
tb(s1, Inches(1.5), Inches(1.3), Inches(8.5), Inches(0.8),
   "业内率先实现上海地区\n安防专网全覆盖", fs=34, color=BLUE, bold=True)
tb(s1, Inches(1.5), Inches(2.4), Inches(8.5), Inches(0.5),
   "实现监控中心的多场所", fs=20, color=SILVER)

# 左侧三个亮点 —— 横向卡片
items = [
    ("01", "安防专网", "上海地区全覆盖\n业内率先完成部署\n统一管控 · 集中值守"),
    ("02", "多场所监控中心", "打破物理空间限制\n远程集中值守\n降低人力成本"),
    ("03", "核心价值", "统一管控体系\n提升响应速度\n安全事件秒级处置"),
]
for i, (num, title, desc) in enumerate(items):
    y = Inches(3.4) + i * Inches(1.5)
    # 编号圆
    circle(s1, Inches(1.5), y + Inches(0.1), Inches(0.45), BLUE); 
    tb(s1, Inches(1.5), y + Inches(0.1), Inches(0.45), Inches(0.45), num, fs=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # 竖线
    rect(s1, Inches(2.15), y, Inches(0.025), Inches(1.2), BLUE_L)
    # 标题
    tb(s1, Inches(2.45), y - Inches(0.05), Inches(3), Inches(0.4), title, fs=16, color=BLUE, bold=True)
    # 描述
    tb(s1, Inches(2.45), y + Inches(0.35), Inches(3.5), Inches(0.9), desc, fs=12, color=GRAY)

# 右侧 —— 精致的网络拓扑示意
# 外圈虚线环（用8个小圆点模拟）
import math
cx, cy = Inches(11.8), Inches(4.8)
for a in range(0, 360, 30):
    r = Inches(2.3)
    px = cx + r * math.cos(math.radians(a))
    py = cy + r * math.sin(math.radians(a))
    circle(s1, px - Inches(0.04), py - Inches(0.04), Inches(0.08), BLUE_L)

# 内圈环
ring = slide_shapes = s1.shapes.add_shape(MSO_SHAPE.OVAL, cx-Inches(2.0), cy-Inches(2.0), Inches(4.0), Inches(4.0))
ring.fill.background(); ring.line.color.rgb = BLUE_L; ring.line.width = Pt(1)

# 中心
center = circle(s1, cx-Inches(0.6), cy-Inches(0.6), Inches(1.2), BLUE)
tb(s1, cx-Inches(0.6), cy-Inches(0.6), Inches(1.2), Inches(1.2), "监控\n中心", fs=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# 6个卫星节点
for i in range(6):
    a_deg = -90 + i * 60
    a = math.radians(a_deg)
    r = Inches(2.3)
    nx = cx + r * math.cos(a) - Inches(0.23)
    ny = cy + r * math.sin(a) - Inches(0.23)
    node = circle(s1, nx, ny, Inches(0.46), SILVER_L, BLUE, 1.5)
    tb(s1, nx, ny, Inches(0.46), Inches(0.46), f"网点{i+1}", fs=7, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    # 连线（简化为细线）
    line(s1, cx-Inches(0.02)+Inches(0.6), cy, nx+Inches(0.23), ny+Inches(0.23), BLUE_L, 0.5)

# 底部
rect(s1, Inches(0), Inches(8.25), Inches(16), Inches(0.75), LIGHT)
tb(s1, Inches(1.5), Inches(8.3), Inches(13), Inches(0.55),
   "【下阶段】加快推进各地分行安防专网建设", fs=13, color=SILVER)

# page number
tb(s1, Inches(14.8), Inches(8.55), Inches(1), Inches(0.35), "01 / 03", fs=9, color=SILVER, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════
# PAGE 2 — 重点部位24小时监测
# ═══════════════════════════════════════════
s2 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s2, WHITE)
rect(s2, Inches(0), Inches(0), Inches(16), Inches(0.06), BLUE)
rect(s2, Inches(0), Inches(0.06), Inches(16), Inches(0.02), BLUE2)
corner_decor(s2, Inches(0.6), Inches(0.5), Inches(0.4), BLUE)
corner_decor(s2, Inches(15.0), Inches(8.1), Inches(0.4), SILVER)
tb(s2, Inches(11.5), Inches(0.3), Inches(4), Inches(2.5), "02", fs=120, color=BLUE_L, bold=True, align=PP_ALIGN.RIGHT)

# 标题
tb(s2, Inches(1.5), Inches(1.2), Inches(10), Inches(0.9),
   "业内首家对办公场所、重控场所\n重点部位 24 小时监测", fs=32, color=BLUE, bold=True)
tb(s2, Inches(1.5), Inches(2.2), Inches(10), Inches(0.45),
   "发现问题实时通知整改  ·  全流程闭环管理", fs=18, color=SILVER)

# 大图区 —— 三列纵向时间流
card_w = Inches(3.6); card_h = Inches(5.0); cy2 = Inches(3.0); gap = Inches(0.6)
start_x = Inches(1.5)
phases = [
    ("感知层", "物联网传感器网络", 
     ["摄像机实时画面采集", "燃气阀门状态监测", "门禁通行记录抓取", "环境温湿度感知"],
     BLUE),
    ("识别层", "AI视觉智能分析",
     ["异常行为自动识别", "人员闯入边界预警", "设备离线秒级告警", "风险事件分级推送"],
     BLUE2),
    ("处置层", "闭环整改追溯",
     ["多渠道通知到责任人", "线上流转处置工单", "整改结果复核归档", "全流程可追溯可审计"],
     RGBColor(0x00, 0x55, 0xA5)),
]
for i, (label, sub, points, accent) in enumerate(phases):
    x = start_x + i*(card_w + gap)
    # 卡片背景
    rect(s2, x, cy2, card_w, card_h, CARD_BG, SILVER_L, 0.5)
    # 顶部色条
    rect(s2, x, cy2, card_w, Inches(0.06), accent)
    # 层级标签
    rect(s2, x, cy2 + Inches(0.3), card_w, Inches(0.55), accent)
    tb(s2, x, cy2 + Inches(0.3), card_w, Inches(0.55), label, fs=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # 副标题
    tb(s2, x + Inches(0.25), cy2 + Inches(1.1), card_w - Inches(0.5), Inches(0.45), sub, fs=13, color=BLUE, bold=True)
    # 要点列表
    y_offset = cy2 + Inches(1.7)
    for j, pt in enumerate(points):
        # 小圆点
        circle(s2, x + Inches(0.25), y_offset + Inches(0.05), Inches(0.1), accent)
        tb(s2, x + Inches(0.5), y_offset - Inches(0.02), card_w - Inches(0.75), Inches(0.55), pt, fs=11, color=DARK)
        y_offset += Inches(0.6)
    
    # 底部箭头（除最后一个）
    if i < 2:
        tb(s2, x + card_w + Inches(0.05), cy2 + card_h/2 - Inches(0.25), Inches(0.35), Inches(0.5),
           "→", fs=24, color=SILVER, bold=True, align=PP_ALIGN.CENTER)

# 底部
rect(s2, Inches(0), Inches(8.25), Inches(16), Inches(0.75), LIGHT)
tb(s2, Inches(1.5), Inches(8.3), Inches(13), Inches(0.55),
   "【下阶段】加快推进各地分行安防专网建设", fs=13, color=SILVER)
tb(s2, Inches(14.8), Inches(8.55), Inches(1), Inches(0.35), "02 / 03", fs=9, color=SILVER, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════
# PAGE 3 — 网络安全管理平台 6大功能
# ═══════════════════════════════════════════
s3 = prs.slides.add_slide(prs.slide_layouts[6]); bg(s3, WHITE)
rect(s3, Inches(0), Inches(0), Inches(16), Inches(0.06), BLUE)
rect(s3, Inches(0), Inches(0.06), Inches(16), Inches(0.02), BLUE2)
corner_decor(s3, Inches(0.6), Inches(0.5), Inches(0.4), BLUE)
corner_decor(s3, Inches(15.0), Inches(8.1), Inches(0.4), SILVER)
tb(s3, Inches(11.5), Inches(0.3), Inches(4), Inches(2.5), "03", fs=120, color=BLUE_L, bold=True, align=PP_ALIGN.RIGHT)

# 标题
tb(s3, Inches(1.5), Inches(1.2), Inches(10), Inches(0.9),
   "搭建业内首个网络安全管理平台", fs=32, color=BLUE, bold=True)
tb(s3, Inches(1.5), Inches(2.1), Inches(10), Inches(0.45),
   "六大功能监测模型  ·  全维度安全防护体系", fs=18, color=SILVER)

# 中心布局 —— 六边形放射
cx3, cy3 = Inches(8.0), Inches(5.3)

# 中心盾牌
shield = circle(s3, cx3-Inches(1.05), cy3-Inches(1.05), Inches(2.1), BLUE)
# 内白圈
inner = circle(s3, cx3-Inches(0.85), cy3-Inches(0.85), Inches(1.7), WHITE, BLUE, 2)
# 盾牌图标文字
tb(s3, cx3-Inches(0.85), cy3-Inches(0.85), Inches(1.7), Inches(1.7),
   "安防\n网络\n安全", fs=18, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

# 外环装饰
outer_ring = s3.shapes.add_shape(MSO_SHAPE.OVAL, cx3-Inches(2.6), cy3-Inches(2.6), Inches(5.2), Inches(5.2))
outer_ring.fill.background(); outer_ring.line.color.rgb = BLUE_L; outer_ring.line.width = Pt(1)

# 6个功能节点
funcs = [
    ("在线准入", "设备身份验证\n安全接入控制", "🛡"),
    ("防外部攻击", "防火墙策略\n入侵检测防御", "🔒"),
    ("隔离内网侵袭", "微隔离技术\n横向移动防护", "🛡"),
    ("扫描系统漏洞", "脆弱性自动巡检\n风险评级修复", "🔍"),
    ("网络稳定性监测", "流量行为分析\n异常精准定位", "📊"),
    ("入侵预警", "实时威胁告警\n自动响应处置", "⚡"),
]
for i, (title, desc, icon) in enumerate(funcs):
    a = math.radians(-90 + i * 60)
    r = Inches(3.1)
    nx = cx3 + r * math.cos(a) - Inches(0.65)
    ny = cy3 + r * math.sin(a) - Inches(0.65)
    
    # 节点圆
    node = circle(s3, nx, ny, Inches(1.3), CARD_BG, BLUE, 2)
    # 小色块顶标
    rect(s3, nx + Inches(0.15), ny + Inches(0.1), Inches(1.0), Inches(0.04), BLUE)
    # 标题
    tb(s3, nx + Inches(0.1), ny + Inches(0.2), Inches(1.1), Inches(0.35), title, fs=11, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    # 描述
    tb(s3, nx + Inches(0.1), ny + Inches(0.55), Inches(1.1), Inches(0.55), desc, fs=8, color=GRAY, align=PP_ALIGN.CENTER)
    
    # 连线从中心到节点
    sx = cx3 + Inches(0.8) * math.cos(a)
    sy = cy3 + Inches(0.8) * math.sin(a)
    line(s3, sx, sy, nx+Inches(0.65), ny+Inches(0.65), BLUE_L, 1.5)

# 底部
rect(s3, Inches(0), Inches(8.25), Inches(16), Inches(0.75), LIGHT)
tb(s3, Inches(1.5), Inches(8.3), Inches(13), Inches(0.55),
   "【下阶段】加快推进各地分行安防专网建设", fs=13, color=SILVER)
tb(s3, Inches(14.8), Inches(8.55), Inches(1), Inches(0.35), "03 / 03", fs=9, color=SILVER, align=PP_ALIGN.RIGHT)

# ── 保存 ──
out = '/root/.openclaw/workspace/projects/security-video/output/镜6_三个要点_PPT_v2.pptx'
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print(f'✅ {out}')

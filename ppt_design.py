#!/usr/bin/env python3
"""
智慧安防体系 — 4张内容页
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BG   = RGBColor(0x0B, 0x0E, 0x14)
CARD = RGBColor(0x13, 0x17, 0x20)
A1   = RGBColor(0x00, 0xD4, 0xFF)   # 蓝
A2   = RGBColor(0x7C, 0x3A, 0xED)   # 紫
A3   = RGBColor(0x10, 0xB9, 0x81)   # 绿
A4   = RGBColor(0xF5, 0x9E, 0x0B)   # 琥珀
GOLD = RGBColor(0xF5, 0xD0, 0x6B)
W    = RGBColor(0xFF, 0xFF, 0xFF)
G1   = RGBColor(0x94, 0xA3, 0xB8)
G2   = RGBColor(0x64, 0x74, 0x8B)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

def bg(s): s.background.fill.solid(); s.background.fill.fore_color.rgb = BG

def r(s, l, t, w, h, fill=None, rad=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rad else MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.line.fill.background()
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else: sh.fill.background()
    return sh

def ln(s, x1, y1, x2, y2, c=A1, w=Pt(1)):
    co = s.shapes.add_connector(1, x1, y1, x2, y2)
    co.line.color.rgb = c; co.line.width = w

def ci(s, l, t, d, fill=None, ol=None):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, l, t, d, d)
    sh.line.fill.background()
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if ol: sh.line.color.rgb = ol; sh.line.width = Pt(1.5)

def tb(s, l, t, w, h, txt, sz=Pt(14), c=W, b=False, a=PP_ALIGN.LEFT):
    bx = s.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt; p.font.size = sz; p.font.color.rgb = c; p.font.bold = b
    p.font.name = '微软雅黑'; p.alignment = a
    return bx

def mt(s, l, t, w, h, lines):
    bx = s.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    for i, (txt, sz, c, b, a) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt; p.font.size = sz; p.font.color.rgb = c; p.font.bold = b
        p.font.name = '微软雅黑'; p.alignment = a; p.space_after = Pt(4)

# ═══════════════════════════════════
# P1 — 预警信息看得清
# ═══════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
r(s, Inches(0), Inches(0), Inches(0.15), Inches(7.5), fill=A2)
# 标题
tb(s, Inches(0.6), Inches(0.35), Inches(2.5), Inches(0.55), '01', Pt(52), A2, True)
tb(s, Inches(3.2), Inches(0.5), Inches(7), Inches(0.55), '预警信息看得清', Pt(36), W, True)
r(s, Inches(3.2), Inches(1.2), Inches(1.5), Pt(3), fill=A2)
tb(s, Inches(3.2), Inches(1.4), Inches(5), Inches(0.3), 'EARLY WARNING VISIBILITY', Pt(12), A2)

# 主卡片
r(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.8), CARD, rad=Pt(10))
r(s, Inches(0.6), Inches(2.0), Inches(6.0), Pt(4), fill=A2)

mt(s, Inches(1.2), Inches(2.5), Inches(5), Inches(3.8), [
    ('消防云平台', Pt(30), A2, True, PP_ALIGN.LEFT),
    ('', Pt(10), W, False, PP_ALIGN.LEFT),
    ('🔥 消防设施感知', Pt(18), W, True, PP_ALIGN.LEFT),
    ('实时监测烟雾 · 温度 · 火焰传感器', Pt(14), G1, False, PP_ALIGN.LEFT),
    ('IoT设备全链路接入，秒级数据采集与回传', Pt(14), G1, False, PP_ALIGN.LEFT),
    ('', Pt(14), W, False, PP_ALIGN.LEFT),
    ('📹 视频联动', Pt(18), W, True, PP_ALIGN.LEFT),
    ('预警触发自动调取周边摄像头', Pt(14), G1, False, PP_ALIGN.LEFT),
    ('AI图像识别辅助判断火情，精准定位', Pt(14), G1, False, PP_ALIGN.LEFT),
    ('', Pt(14), W, False, PP_ALIGN.LEFT),
    ('IoT感知  ·  AI识别  ·  自动联动  ·  秒级预警', Pt(15), A1, True, PP_ALIGN.LEFT),
])

# 右侧链路图
r(s, Inches(7.2), Inches(2.0), Inches(5.5), Inches(4.8), CARD, rad=Pt(10))
tb(s, Inches(8.3), Inches(2.5), Inches(3.5), Inches(0.5), '⚡ 响应链路', Pt(20), W, True)

steps = ['传感器感知', '数据汇聚', 'AI研判', '预警触发', '视频联动', '处置闭环']
cols = [A2, A2, A1, GOLD, A4, A3]
for i, (txt, cl) in enumerate(zip(steps, cols)):
    x = Inches(7.6 + i * 0.9)
    ci(s, x, Inches(3.7), Inches(0.7), fill=cl)
    tb(s, x, Inches(3.8), Inches(0.7), Inches(0.5), txt, Pt(9), W, True, PP_ALIGN.CENTER)
    if i < 5:
        ln(s, x + Inches(0.72), Inches(4.05), x + Inches(0.88), Inches(4.05), A1, Pt(1.5))

for i, (v, lb) in enumerate([('＜3s', '预警延迟'), ('99.9%', '感知覆盖率'), ('7×24', '不间断')]):
    x = Inches(7.8 + i * 1.85)
    tb(s, x, Inches(5.0), Inches(1.6), Inches(0.45), v, Pt(24), A1, True, PP_ALIGN.CENTER)
    tb(s, x, Inches(5.5), Inches(1.6), Inches(0.3), lb, Pt(11), G1, False, PP_ALIGN.CENTER)

# 底部碎线
ln(s, Inches(0.6), Inches(7.1), Inches(3), Inches(7.1), A2, Pt(1))

# ═══════════════════════════════════
# P2 — 管理手势看得清
# ═══════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
r(s, Inches(0), Inches(0), Inches(0.15), Inches(7.5), fill=A3)
tb(s, Inches(0.6), Inches(0.35), Inches(2.5), Inches(0.55), '02', Pt(52), A3, True)
tb(s, Inches(3.2), Inches(0.5), Inches(7), Inches(0.55), '管理手势看得清', Pt(36), W, True)
r(s, Inches(3.2), Inches(1.2), Inches(1.5), Pt(3), fill=A3)
tb(s, Inches(3.2), Inches(1.4), Inches(5), Inches(0.3), 'MANAGEMENT VISIBILITY', Pt(12), A3)

r(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.8), CARD, rad=Pt(10))
r(s, Inches(0.6), Inches(2.0), Inches(6.0), Pt(4), fill=A3)

mt(s, Inches(1.2), Inches(2.5), Inches(5), Inches(3.8), [
    ('智慧楼宇管理平台', Pt(30), A3, True, PP_ALIGN.LEFT),
    ('', Pt(10), W, False, PP_ALIGN.LEFT),
    ('🏢 总行本级办公楼宇可视化管理', Pt(18), W, True, PP_ALIGN.LEFT),
    ('BIM + 3D数字孪生，楼宇结构 / 设备 / 管线一屏掌控', Pt(14), G1, False, PP_ALIGN.LEFT),
    ('', Pt(14), W, False, PP_ALIGN.LEFT),
    ('📋 核心能力', Pt(16), W, True, PP_ALIGN.LEFT),
    ('空间管理  ·  能耗监控  ·  设备运维  ·  访客轨迹', Pt(14), G1, False, PP_ALIGN.LEFT),
    ('安防态势  ·  环境监测  ·  应急指挥  ·  资产盘点', Pt(14), G1, False, PP_ALIGN.LEFT),
    ('', Pt(14), W, False, PP_ALIGN.LEFT),
    ('一屏统览 · 一键调度 · 一图决策', Pt(16), A3, True, PP_ALIGN.LEFT),
])

r(s, Inches(7.2), Inches(2.0), Inches(5.5), Inches(4.8), CARD, rad=Pt(10))
tb(s, Inches(8.3), Inches(2.5), Inches(3.5), Inches(0.5), '🏗️ 数字化楼宇', Pt(20), W, True)

for i in range(1, 9):
    y = Inches(3.0 + (8 - i) * 0.42)
    w = Inches(1.5 + i * 0.55)
    r(s, Inches(8.6), y, w, Inches(0.37), fill=A3)
    tb(s, Inches(7.6), y, Inches(0.6), Inches(0.37), f'F{i}', Pt(9), G1, False, PP_ALIGN.RIGHT)

tb(s, Inches(7.5), Inches(6.3), Inches(5.2), Inches(0.3), 'BIM数字孪生  ·  IoT全感知  ·  3D可视化', Pt(13), A3, True, PP_ALIGN.CENTER)
ln(s, Inches(0.6), Inches(7.1), Inches(3), Inches(7.1), A3, Pt(1))

# ═══════════════════════════════════
# P3 — 处置流程看得清
# ═══════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
r(s, Inches(0), Inches(0), Inches(0.15), Inches(7.5), fill=GOLD)
tb(s, Inches(0.6), Inches(0.35), Inches(2.5), Inches(0.55), '03', Pt(52), GOLD, True)
tb(s, Inches(3.2), Inches(0.5), Inches(7), Inches(0.55), '处置流程看得清', Pt(36), W, True)
r(s, Inches(3.2), Inches(1.2), Inches(1.5), Pt(3), fill=GOLD)
tb(s, Inches(3.2), Inches(1.4), Inches(5), Inches(0.3), 'PROCESS VISIBILITY', Pt(12), GOLD)

r(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.8), CARD, rad=Pt(10))
r(s, Inches(0.6), Inches(2.0), Inches(6.0), Pt(4), fill=GOLD)

mt(s, Inches(1.2), Inches(2.5), Inches(5), Inches(3.8), [
    ('金库智能安防系统', Pt(30), GOLD, True, PP_ALIGN.LEFT),
    ('', Pt(10), W, False, PP_ALIGN.LEFT),
    ('🔐 金库安全数字化统一管理', Pt(18), W, True, PP_ALIGN.LEFT),
    ('多因子认证 · 行为分析 · 全流程可追溯', Pt(14), G1, False, PP_ALIGN.LEFT),
    ('', Pt(14), W, False, PP_ALIGN.LEFT),
    ('🛡️ 安全层级', Pt(16), W, True, PP_ALIGN.LEFT),
    ('第一层  生物识别 + 多因子认证', Pt(14), A2, False, PP_ALIGN.LEFT),
    ('第二层  智能视频分析 + 异常行为检测', Pt(14), A4, False, PP_ALIGN.LEFT),
    ('第三层  全流程审计追溯 + 区块链存证', Pt(14), A1, False, PP_ALIGN.LEFT),
    ('', Pt(14), W, False, PP_ALIGN.LEFT),
    ('统一管控 · 实时告警 · 合规审计 · 零信任架构', Pt(16), GOLD, True, PP_ALIGN.LEFT),
])

r(s, Inches(7.2), Inches(2.0), Inches(5.5), Inches(4.8), CARD, rad=Pt(10))
tb(s, Inches(8.3), Inches(2.5), Inches(3.5), Inches(0.5), '🛡️ 安全防护环', Pt(20), W, True)

for ri, (rad, cl, lb, vo) in enumerate([
    (Inches(1.5), A2, '身份认证', Inches(0.45)),
    (Inches(1.05), A4, '行为分析', Inches(0.30)),
    (Inches(0.6), GOLD, '审计追溯', Inches(0.15)),
]):
    ci(s, Inches(9.4) - rad, Inches(3.7) - rad, rad * 2, ol=cl)
    tb(s, Inches(9.4) - Inches(0.5), Inches(3.7) + rad - vo,
        Inches(1.0), Inches(0.25), lb, Pt(9), cl, True, PP_ALIGN.CENTER)

ci(s, Inches(9.13), Inches(3.43), Inches(0.55), fill=GOLD)
tb(s, Inches(9.13), Inches(3.5), Inches(0.55), Inches(0.3), '核心', Pt(9), BG, True, PP_ALIGN.CENTER)

tb(s, Inches(7.5), Inches(6.3), Inches(5.2), Inches(0.3), '全链路数字化  ·  零信任架构  ·  实时风控', Pt(13), GOLD, True, PP_ALIGN.CENTER)
ln(s, Inches(0.6), Inches(7.1), Inches(3), Inches(7.1), GOLD, Pt(1))

# ═══════════════════════════════════
# P4 — 下阶段
# ═══════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
r(s, Inches(0), Inches(0), Inches(0.15), Inches(7.5), fill=A1)
tb(s, Inches(0.6), Inches(0.35), Inches(2.5), Inches(0.55), 'NEXT', Pt(52), A1, True)
tb(s, Inches(3.2), Inches(0.5), Inches(7), Inches(0.55), '下阶段规划', Pt(36), W, True)
r(s, Inches(3.2), Inches(1.2), Inches(1.5), Pt(3), fill=A1)
tb(s, Inches(3.2), Inches(1.4), Inches(5), Inches(0.3), 'NEXT PHASE', Pt(12), A1)

# 四张卡片
items = [
    ('🧊', '搭建数据湖', '统一数据底座\n汇聚消防 · 楼宇 · 安防全维度数据', A1),
    ('🔗', '打通四大平台\n数据链路', '消防云 · 楼宇管理 · 金库安防\n· E事通数据互联互通', A3),
    ('📱', '安防APP\n接入E事通', '移动端统一入口\n随时随地掌握安防态势', A4),
    ('⚠️', '构建统一风险\n预警视图', '一张图呈现全行风险态势\n智能分级预警 · 联动处置', GOLD),
]

for i, (icon, title, desc, cl) in enumerate(items):
    x = Inches(0.6 + i * 3.15)
    y = Inches(2.1)
    cw = Inches(2.85)
    ch = Inches(3.4)
    r(s, x, y, cw, ch, CARD, rad=Pt(10))
    r(s, x, y, cw, Pt(4), fill=cl)
    tb(s, x + Inches(0.25), y + Inches(0.3), Inches(2.3), Inches(0.45), icon, Pt(30), W)
    tb(s, x + Inches(0.25), y + Inches(1.2), Inches(2.35), Inches(0.9), title, Pt(18), W, True)
    tb(s, x + Inches(0.25), y + Inches(2.3), Inches(2.35), Inches(0.9), desc, Pt(11), G1)

# 时间轴
ln(s, Inches(0.6), Inches(6.0), Inches(12.7), Inches(6.0), LINE_C := RGBColor(0x1E,0x29,0x3B), Pt(1.5))
for i, (t, lb) in enumerate([('Q3 2026', '数据湖搭建'), ('Q4 2026', '平台对接'), ('Q1 2027', 'APP集成'), ('Q2 2027', '全景上线')]):
    x = Inches(0.6 + i * 3.15)
    ci(s, x + Inches(1.1), Inches(5.85), Inches(0.28), fill=A1)
    tb(s, x, Inches(6.2), Inches(2.45), Inches(0.25), t, Pt(11), W, True, PP_ALIGN.CENTER)
    tb(s, x, Inches(6.5), Inches(2.45), Inches(0.25), lb, Pt(10), G1, False, PP_ALIGN.CENTER)

ln(s, Inches(0.6), Inches(7.1), Inches(3), Inches(7.1), A1, Pt(1))

# 保存
out = '/root/.openclaw/workspace/智慧安防体系_PPT.pptx'
prs.save(out)
print(f'✅ {len(prs.slides)} pages → {out}')
